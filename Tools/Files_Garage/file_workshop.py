"""
File Workshop v4 — Universal Local File Toolkit
================================================
CONVERT:   PDF · DOCX · XLSX · PPTX · CSV · TXT · HTML · Images · Audio · Video
SPLIT:     Each page | Range | Custom groups  (PDF / PPTX)
MERGE:     Combine PDFs / images / PPTX slides into one file
ORGANISE:  Resequence · Delete · Rotate · Reverse  (PDF / PPTX)
STAMP:     Watermark text / PDF overlay
PROTECT:   Encrypt / decrypt PDF with password
COMPRESS:  Reduce PDF file size
METADATA:  View and edit PDF / Office metadata

Requirements: see requirements.txt
"""

import os, sys, re, csv, threading, subprocess, shutil, tempfile
from pathlib import Path
from typing import List, Dict, Optional
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext

# ═══════════════════════════════════════════════════════════════════════════════
# Dependency detection
# ═══════════════════════════════════════════════════════════════════════════════

def _try(pkg):
    try: __import__(pkg); return True
    except ImportError: return False

HAS_PYPDF      = _try("pypdf")
HAS_PDFPLUMBER = _try("pdfplumber")
HAS_DOCX       = _try("docx")
HAS_PIL        = _try("PIL")
HAS_PDF2IMAGE  = _try("pdf2image")
HAS_REPORTLAB  = _try("reportlab")
HAS_WEASYPRINT = _try("weasyprint")
HAS_OPENPYXL   = _try("openpyxl")
HAS_PPTX       = _try("pptx")
HAS_PANDAS     = _try("pandas")
HAS_FFMPEG     = shutil.which("ffmpeg") is not None
HAS_LIBREOFFICE = bool(shutil.which("libreoffice") or shutil.which("soffice"))

if HAS_PYPDF:
    from pypdf import PdfReader, PdfWriter
if HAS_PDFPLUMBER:
    import pdfplumber
if HAS_DOCX:
    from docx import Document
    from docx.shared import Pt, Inches
if HAS_PIL:
    from PIL import Image, ImageDraw, ImageFont
if HAS_PDF2IMAGE:
    from pdf2image import convert_from_path
if HAS_OPENPYXL:
    import openpyxl
if HAS_PPTX:
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt
if HAS_PANDAS:
    import pandas as pd

# ═══════════════════════════════════════════════════════════════════════════════
# File categorisation
# ═══════════════════════════════════════════════════════════════════════════════

IMAGE_EXTS = {".png",".jpg",".jpeg",".webp",".bmp",".gif",".tiff",".tif",".ico"}
AUDIO_EXTS = {".mp3",".wav",".ogg",".flac",".aac",".m4a",".wma"}
VIDEO_EXTS = {".mp4",".avi",".mov",".mkv",".webm",".flv",".wmv",".m4v"}
EXCEL_EXTS = {".xlsx",".xls",".xlsm",".ods"}
PPTX_EXTS  = {".pptx",".ppt",".odp"}
CSV_EXTS   = {".csv",".tsv"}

def cat(path: str) -> str:
    e = Path(path).suffix.lower()
    if e in IMAGE_EXTS: return "image"
    if e in AUDIO_EXTS: return "audio"
    if e in VIDEO_EXTS: return "video"
    if e in EXCEL_EXTS: return "excel"
    if e in PPTX_EXTS:  return "pptx"
    if e in CSV_EXTS:   return "csv"
    if e == ".pdf":     return "pdf"
    if e == ".docx":    return "docx"
    if e == ".txt":     return "txt"
    if e in {".html",".htm"}: return "html"
    return "unknown"

def cat_icon(c: str) -> str:
    return {"pdf":"📄","image":"🖼","audio":"🎵","video":"🎬",
            "docx":"📝","txt":"📃","html":"🌐","excel":"📊",
            "pptx":"📽","csv":"📋"}.get(c, "📎")

# ═══════════════════════════════════════════════════════════════════════════════
# Page / row string parsers
# ═══════════════════════════════════════════════════════════════════════════════

def parse_pages(s: str, total: int) -> List[int]:
    pages = set()
    for part in s.split(","):
        part = part.strip()
        if "-" in part:
            a, b = part.split("-", 1)
            try: pages.update(range(int(a), int(b)+1))
            except: pass
        else:
            try: pages.add(int(part))
            except: pass
    return sorted(p for p in pages if 1 <= p <= total)

def parse_groups(s: str, total: int) -> List[List[int]]:
    return [g for g in (parse_pages(c.strip(), total) for c in s.split("|")) if g]

# ═══════════════════════════════════════════════════════════════════════════════
# LibreOffice universal helper  (XLSX/PPTX/DOCX → PDF via headless LO)
# ═══════════════════════════════════════════════════════════════════════════════

def lo_convert(src: str, dst_fmt: str, out_dir: str) -> str:
    """Use LibreOffice headless to convert src to dst_fmt, save in out_dir."""
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if not lo:
        raise RuntimeError(
            "LibreOffice not found.\n"
            "  Ubuntu/Debian: sudo apt install libreoffice\n"
            "  macOS:         brew install --cask libreoffice\n"
            "  Windows:       https://libreoffice.org")
    os.makedirs(out_dir, exist_ok=True)
    r = subprocess.run(
        [lo, "--headless", "--convert-to", dst_fmt, "--outdir", out_dir, src],
        capture_output=True, text=True)
    stem = Path(src).stem
    expected = Path(out_dir) / f"{stem}.{dst_fmt}"
    if not expected.exists():
        raise RuntimeError(f"LibreOffice conversion failed:\n{r.stderr[-600:]}")
    return str(expected)

# ═══════════════════════════════════════════════════════════════════════════════
# ── PDF conversions ───────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def pdf_page_count(src): return len(PdfReader(src).pages)

def pdf_to_txt(src, dst, pages=None):
    if not HAS_PDFPLUMBER: raise ImportError("pip install pdfplumber")
    with pdfplumber.open(src) as pdf:
        total = len(pdf.pages)
        target = pages or list(range(1, total+1))
        lines = []
        for p in target:
            if 1 <= p <= total:
                lines.append(f"\n{'='*60}\nPAGE {p}\n{'='*60}\n")
                lines.append(pdf.pages[p-1].extract_text() or "")
    Path(dst).write_text("\n".join(lines), encoding="utf-8")
    return dst

def pdf_to_docx(src, dst, pages=None):
    if not (HAS_PDFPLUMBER and HAS_DOCX): raise ImportError("pip install pdfplumber python-docx")
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    with pdfplumber.open(src) as pdf:
        total = len(pdf.pages)
        target = pages or list(range(1, total+1))
        for p in target:
            if 1 <= p <= total:
                doc.add_heading(f"Page {p}", 2)
                text = pdf.pages[p-1].extract_text() or ""
                for para in text.split("\n\n"):
                    if para.strip(): doc.add_paragraph(para.strip())
                doc.add_page_break()
    doc.save(dst); return dst

def pdf_to_images(src, dst_dir, fmt="png", pages=None, dpi=150):
    if not HAS_PDF2IMAGE: raise ImportError("pip install pdf2image  (+poppler)")
    os.makedirs(dst_dir, exist_ok=True)
    kw = {"dpi": dpi, "fmt": fmt}
    if pages: kw["first_page"], kw["last_page"] = min(pages), max(pages)
    imgs = convert_from_path(src, **kw)
    out = []
    for i, img in enumerate(imgs):
        pg = pages[i] if pages and i < len(pages) else i+1
        p = os.path.join(dst_dir, f"page_{pg}.{fmt}")
        img.save(p); out.append(p)
    return out

def pdf_to_html(src, dst, pages=None):
    if not HAS_PDFPLUMBER: raise ImportError("pip install pdfplumber")
    with pdfplumber.open(src) as pdf:
        total = len(pdf.pages)
        target = pages or list(range(1, total+1))
        parts = [f"<html><head><meta charset='utf-8'><title>{Path(src).stem}</title>"
                 f"<style>body{{font-family:Georgia,serif;max-width:860px;margin:auto;padding:2em}}"
                 f"h2{{border-bottom:2px solid #c00}}</style></head><body>"]
        for p in target:
            if 1 <= p <= total:
                text = pdf.pages[p-1].extract_text() or ""
                safe = text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                parts.append(f"<h2>Page {p}</h2><pre style='white-space:pre-wrap'>{safe}</pre><hr>")
        parts.append("</body></html>")
    Path(dst).write_text("\n".join(parts), encoding="utf-8"); return dst

def pdf_to_xlsx(src, dst):
    """Extract tables from PDF into Excel workbook (one sheet per page)."""
    if not (HAS_PDFPLUMBER and HAS_OPENPYXL):
        raise ImportError("pip install pdfplumber openpyxl")
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    with pdfplumber.open(src) as pdf:
        for i, page in enumerate(pdf.pages):
            tables = page.extract_tables()
            ws = wb.create_sheet(title=f"Page_{i+1}")
            if tables:
                for table in tables:
                    for row in table:
                        ws.append([c or "" for c in row])
                    ws.append([])  # blank separator
            else:
                # fall back to raw text in column A
                text = page.extract_text() or ""
                for line in text.split("\n"):
                    ws.append([line])
    wb.save(dst); return dst

def pdf_to_pptx(src, dst, dpi=150):
    """Convert each PDF page to an image slide in PowerPoint."""
    if not (HAS_PDF2IMAGE and HAS_PPTX):
        raise ImportError("pip install pdf2image python-pptx  (+poppler)")
    imgs = convert_from_path(src, dpi=dpi)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # completely blank
    for img in imgs:
        slide = prs.slides.add_slide(blank_layout)
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tf:
            tmp = tf.name
        img.save(tmp)
        slide.shapes.add_picture(tmp, 0, 0,
                                  width=prs.slide_width,
                                  height=prs.slide_height)
        os.unlink(tmp)
    prs.save(dst); return dst

# ═══════════════════════════════════════════════════════════════════════════════
# ── Excel conversions ─────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def excel_to_pdf(src, dst):
    """Excel → PDF via LibreOffice."""
    out_dir = str(Path(dst).parent)
    result = lo_convert(src, "pdf", out_dir)
    if result != dst: os.rename(result, dst)
    return dst

def excel_to_csv(src, dst, sheet_index=0):
    """Export first (or specified) sheet of Excel to CSV."""
    if not HAS_OPENPYXL: raise ImportError("pip install openpyxl")
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    sheet = wb.worksheets[sheet_index]
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        for row in sheet.iter_rows(values_only=True):
            writer.writerow([v if v is not None else "" for v in row])
    wb.close(); return dst

def excel_to_txt(src, dst):
    """Export all sheets of Excel to plain text."""
    if not HAS_OPENPYXL: raise ImportError("pip install openpyxl")
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"\n{'='*60}\nSHEET: {ws.title}\n{'='*60}\n")
        for row in ws.iter_rows(values_only=True):
            lines.append("\t".join(str(v) if v is not None else "" for v in row))
    wb.close()
    Path(dst).write_text("\n".join(lines), encoding="utf-8"); return dst

def excel_to_html(src, dst):
    """Convert Excel to HTML table (one section per sheet)."""
    if not HAS_OPENPYXL: raise ImportError("pip install openpyxl")
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    parts = [f"<html><head><meta charset='utf-8'><title>{Path(src).stem}</title>"
             f"<style>body{{font-family:sans-serif;padding:2em}}"
             f"table{{border-collapse:collapse;margin-bottom:2em}}"
             f"td,th{{border:1px solid #ccc;padding:4px 10px;font-size:13px}}"
             f"th{{background:#f0f0f0}}h2{{color:#333}}</style></head><body>"]
    for ws in wb.worksheets:
        parts.append(f"<h2>{ws.title}</h2><table>")
        first_row = True
        for row in ws.iter_rows(values_only=True):
            tag = "th" if first_row else "td"
            cells = "".join(f"<{tag}>{str(v) if v is not None else ''}</{tag}>" for v in row)
            parts.append(f"<tr>{cells}</tr>")
            first_row = False
        parts.append("</table>")
    parts.append("</body></html>")
    wb.close()
    Path(dst).write_text("\n".join(parts), encoding="utf-8"); return dst

def excel_to_docx(src, dst):
    """Export Excel content to a Word document with tables."""
    if not (HAS_OPENPYXL and HAS_DOCX):
        raise ImportError("pip install openpyxl python-docx")
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    for ws in wb.worksheets:
        doc.add_heading(ws.title, 2)
        rows = list(ws.iter_rows(values_only=True))
        if not rows: continue
        ncols = max(len(r) for r in rows)
        table = doc.add_table(rows=len(rows), cols=ncols)
        table.style = "Table Grid"
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = table.cell(ri, ci)
                cell.text = str(val) if val is not None else ""
                if ri == 0:
                    for run in cell.paragraphs[0].runs:
                        run.bold = True
        doc.add_paragraph()
    wb.close(); doc.save(dst); return dst

def csv_to_xlsx(src, dst):
    """CSV → Excel workbook."""
    if not HAS_OPENPYXL: raise ImportError("pip install openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    delim = "\t" if src.lower().endswith(".tsv") else ","
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delim)
        for row in reader: ws.append(row)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    wb.save(dst); return dst

def csv_to_pdf(src, dst):
    """CSV → PDF table via reportlab."""
    if not HAS_REPORTLAB: raise ImportError("pip install reportlab")
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import mm
    rows = []
    delim = "\t" if src.lower().endswith(".tsv") else ","
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delim)
        for row in reader: rows.append(row)
    if not rows: raise ValueError("CSV file is empty")
    doc = SimpleDocTemplate(dst, pagesize=landscape(A4),
                            leftMargin=10*mm, rightMargin=10*mm,
                            topMargin=10*mm, bottomMargin=10*mm)
    styles = getSampleStyleSheet()
    story = [Paragraph(Path(src).stem, styles["Title"]), Spacer(1, 6)]
    ncols = max(len(r) for r in rows)
    padded = [r + [""]*(ncols-len(r)) for r in rows]
    t = Table(padded, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#2e2e4a")),
        ("TEXTCOLOR",  (0,0), (-1,0), colors.white),
        ("FONTNAME",   (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE",   (0,0), (-1,-1), 8),
        ("GRID",       (0,0), (-1,-1), 0.5, colors.grey),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#f5f5f5")]),
        ("VALIGN",     (0,0), (-1,-1), "MIDDLE"),
        ("PADDING",    (0,0), (-1,-1), 4),
    ]))
    story.append(t)
    doc.build(story); return dst

def csv_to_html(src, dst):
    """CSV → HTML table."""
    rows = []
    delim = "\t" if src.lower().endswith(".tsv") else ","
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delim)
        for row in reader: rows.append(row)
    parts = [f"<html><head><meta charset='utf-8'><title>{Path(src).stem}</title>"
             f"<style>body{{font-family:sans-serif;padding:2em}}"
             f"table{{border-collapse:collapse}}td,th{{border:1px solid #ccc;padding:5px 12px;"
             f"font-size:13px}}th{{background:#e8e8f0}}</style></head><body>"
             f"<h2>{Path(src).stem}</h2><table>"]
    for i, row in enumerate(rows):
        tag = "th" if i == 0 else "td"
        cells = "".join(f"<{tag}>{str(v).replace('&','&amp;').replace('<','&lt;')}</{tag}>" for v in row)
        parts.append(f"<tr>{cells}</tr>")
    parts.append("</table></body></html>")
    Path(dst).write_text("\n".join(parts), encoding="utf-8"); return dst

def csv_to_txt(src, dst):
    shutil.copy(src, dst); return dst  # CSV is already text, just copy

def csv_to_docx(src, dst):
    """CSV → Word document with table."""
    if not HAS_DOCX: raise ImportError("pip install python-docx")
    rows = []
    delim = "\t" if src.lower().endswith(".tsv") else ","
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delim)
        for row in reader: rows.append(row)
    if not rows: raise ValueError("CSV is empty")
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    ncols = max(len(r) for r in rows)
    padded = [r + [""]*(ncols-len(r)) for r in rows]
    table = doc.add_table(rows=len(padded), cols=ncols)
    table.style = "Table Grid"
    for ri, row in enumerate(padded):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            if ri == 0:
                for run in cell.paragraphs[0].runs:
                    run.bold = True
    doc.save(dst); return dst

# ═══════════════════════════════════════════════════════════════════════════════
# ── PowerPoint conversions ────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def pptx_to_pdf(src, dst):
    """PPTX → PDF via LibreOffice."""
    out_dir = str(Path(dst).parent)
    result = lo_convert(src, "pdf", out_dir)
    if str(result) != str(dst):
        shutil.move(result, dst)
    return dst

def pptx_to_images(src, dst_dir, fmt="png", dpi=150):
    """PPTX → PNG/JPG by converting to PDF first then rasterising."""
    if not HAS_PDF2IMAGE: raise ImportError("pip install pdf2image  (+poppler)")
    os.makedirs(dst_dir, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        tmp_pdf = os.path.join(tmp, Path(src).stem + ".pdf")
        pptx_to_pdf(src, tmp_pdf)
        imgs = convert_from_path(tmp_pdf, dpi=dpi, fmt=fmt)
        out = []
        for i, img in enumerate(imgs):
            p = os.path.join(dst_dir, f"slide_{i+1}.{fmt}")
            img.save(p); out.append(p)
    return out

def pptx_to_txt(src, dst):
    """Extract all text from PPTX slides."""
    if not HAS_PPTX: raise ImportError("pip install python-pptx")
    prs = Presentation(src)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n{'='*60}\nSLIDE {i}\n{'='*60}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text: lines.append(text)
    Path(dst).write_text("\n".join(lines), encoding="utf-8"); return dst

def pptx_to_html(src, dst):
    """PPTX → HTML (text-based, one section per slide)."""
    if not HAS_PPTX: raise ImportError("pip install python-pptx")
    prs = Presentation(src)
    parts = [f"<html><head><meta charset='utf-8'><title>{Path(src).stem}</title>"
             f"<style>body{{font-family:Georgia,serif;max-width:900px;margin:auto;padding:2em}}"
             f".slide{{border:1px solid #ddd;margin-bottom:2em;padding:1.5em;border-radius:4px}}"
             f"h2{{color:#444;border-bottom:2px solid #7c6af7;padding-bottom:.4em}}"
             f"</style></head><body>"
             f"<h1>{Path(src).stem}</h1>"]
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f'<div class="slide"><h2>Slide {i}</h2>')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        safe = text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                        parts.append(f"<p>{safe}</p>")
        parts.append("</div>")
    parts.append("</body></html>")
    Path(dst).write_text("\n".join(parts), encoding="utf-8"); return dst

def pptx_to_docx(src, dst):
    """PPTX → DOCX (text content preserved per slide)."""
    if not (HAS_PPTX and HAS_DOCX): raise ImportError("pip install python-pptx python-docx")
    prs = Presentation(src)
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    for i, slide in enumerate(prs.slides, 1):
        doc.add_heading(f"Slide {i}", 2)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text: doc.add_paragraph(text)
        doc.add_page_break()
    doc.save(dst); return dst

def pptx_slide_count(src):
    if not HAS_PPTX: return 0
    return len(Presentation(src).slides)

def pptx_split_slides(src, out_dir, prefix="slide"):
    """Split each slide of a PPTX into its own PPTX file."""
    if not HAS_PPTX: raise ImportError("pip install python-pptx")
    os.makedirs(out_dir, exist_ok=True)
    prs = Presentation(src)
    out = []
    for i, slide in enumerate(prs.slides):
        new_prs = Presentation()
        # Match slide dimensions
        new_prs.slide_width  = prs.slide_width
        new_prs.slide_height = prs.slide_height
        layout = new_prs.slide_layouts[6]  # blank
        new_slide = new_prs.slides.add_slide(layout)
        # Copy shapes via XML
        from pptx.oxml.ns import qn
        import copy
        for shape in slide.shapes:
            el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert(2, el)
        p = os.path.join(out_dir, f"{prefix}_{i+1}.pptx")
        new_prs.save(p); out.append(p)
    return out

def pptx_merge(src_list, dst):
    """Merge multiple PPTX files into one, preserving all slides."""
    if not HAS_PPTX: raise ImportError("pip install python-pptx")
    import copy
    from lxml import etree
    base_prs = Presentation(src_list[0])
    for src in src_list[1:]:
        src_prs = Presentation(src)
        for slide in src_prs.slides:
            layout = base_prs.slide_layouts[6]
            new_slide = base_prs.slides.add_slide(layout)
            new_slide.slide_width  = src_prs.slide_width
            new_slide.slide_height = src_prs.slide_height
            for shape in slide.shapes:
                el = copy.deepcopy(shape.element)
                new_slide.shapes._spTree.insert(2, el)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    base_prs.save(dst); return dst

# ═══════════════════════════════════════════════════════════════════════════════
# ── DOCX conversions ──────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def docx_to_txt(src, dst):
    if not HAS_DOCX: raise ImportError("pip install python-docx")
    doc = Document(src)
    Path(dst).write_text("\n".join(p.text for p in doc.paragraphs), encoding="utf-8")
    return dst

def docx_to_html(src, dst):
    if not HAS_DOCX: raise ImportError("pip install python-docx")
    doc = Document(src)
    parts = [f"<html><head><meta charset='utf-8'><title>{Path(src).stem}</title>"
             f"<style>body{{font-family:Georgia,serif;max-width:860px;margin:auto;padding:2em}}</style></head><body>"]
    for p in doc.paragraphs:
        if p.style.name.startswith("Heading"):
            lvl = p.style.name[-1] if p.style.name[-1].isdigit() else "2"
            parts.append(f"<h{lvl}>{p.text}</h{lvl}>")
        else:
            safe = p.text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            parts.append(f"<p>{safe}</p>")
    parts.append("</body></html>")
    Path(dst).write_text("\n".join(parts), encoding="utf-8"); return dst

def docx_to_pdf(src, dst):
    out_dir = str(Path(dst).parent)
    result = lo_convert(src, "pdf", out_dir)
    if str(result) != str(dst): shutil.move(result, dst)
    return dst

# ═══════════════════════════════════════════════════════════════════════════════
# ── TXT conversions ───────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def txt_to_pdf(src, dst):
    if not HAS_REPORTLAB: raise ImportError("pip install reportlab")
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import mm
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    doc = SimpleDocTemplate(dst, pagesize=A4,
                            leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=20*mm, bottomMargin=20*mm)
    styles = getSampleStyleSheet()
    story = []
    for line in text.split("\n"):
        safe = line.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;") or "&nbsp;"
        story.append(Paragraph(safe, styles["Normal"]))
        story.append(Spacer(1, 2))
    doc.build(story); return dst

def txt_to_docx(src, dst):
    if not HAS_DOCX: raise ImportError("pip install python-docx")
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    for line in Path(src).read_text(encoding="utf-8", errors="replace").split("\n"):
        doc.add_paragraph(line)
    doc.save(dst); return dst

def txt_to_html(src, dst):
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    safe = text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
    Path(dst).write_text(
        f"<html><head><meta charset='utf-8'><title>{Path(src).stem}</title>"
        f"<style>body{{font-family:monospace;max-width:900px;margin:auto;padding:2em}}</style></head>"
        f"<body><pre>{safe}</pre></body></html>", encoding="utf-8")
    return dst

# ═══════════════════════════════════════════════════════════════════════════════
# ── HTML conversions ──────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def html_to_txt(src, dst):
    html = Path(src).read_text(encoding="utf-8", errors="replace")
    text = re.sub(r"<[^>]+>", "", html)
    for ent, rep in [("&nbsp;"," "),("&amp;","&"),("&lt;","<"),("&gt;",">")]:
        text = text.replace(ent, rep)
    Path(dst).write_text(text.strip(), encoding="utf-8"); return dst

def html_to_pdf(src, dst):
    if HAS_WEASYPRINT:
        from weasyprint import HTML
        HTML(filename=src).write_pdf(dst); return dst
    wk = shutil.which("wkhtmltopdf")
    if wk:
        r = subprocess.run([wk, src, dst], capture_output=True)
        if Path(dst).exists(): return dst
    raise ImportError("HTML→PDF needs:  pip install weasyprint")

# ═══════════════════════════════════════════════════════════════════════════════
# ── Image conversions ─────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def image_convert(src, dst):
    if not HAS_PIL: raise ImportError("pip install Pillow")
    img = Image.open(src)
    ext = Path(dst).suffix.lower()
    if ext in (".jpg",".jpeg",".bmp") and img.mode in ("RGBA","P","LA"):
        img = img.convert("RGB")
    if ext == ".ico": img = img.resize((256,256), Image.LANCZOS)
    img.save(dst); return dst

def images_to_pdf(src_list, dst):
    if not HAS_PIL: raise ImportError("pip install Pillow")
    imgs = [Image.open(p).convert("RGB") for p in src_list]
    if imgs: imgs[0].save(dst, save_all=True, append_images=imgs[1:])
    return dst

# ═══════════════════════════════════════════════════════════════════════════════
# ── Audio / Video ─────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def ffmpeg_convert(src, dst, extra=None):
    if not HAS_FFMPEG:
        raise RuntimeError("ffmpeg not found.\n  Ubuntu: sudo apt install ffmpeg\n  Mac: brew install ffmpeg")
    cmd = ["ffmpeg","-y","-i",src] + (extra or []) + [dst]
    r = subprocess.run(cmd, capture_output=True, text=True)
    if not Path(dst).exists():
        raise RuntimeError(f"ffmpeg failed:\n{r.stderr[-600:]}")
    return dst

# ═══════════════════════════════════════════════════════════════════════════════
# ── PDF manipulation ──────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def split_each(src, out_dir, prefix="page"):
    os.makedirs(out_dir, exist_ok=True)
    reader = PdfReader(src); out = []
    for i, page in enumerate(reader.pages):
        w = PdfWriter(); w.add_page(page)
        p = os.path.join(out_dir, f"{prefix}_{i+1}.pdf")
        with open(p,"wb") as f: w.write(f)
        out.append(p)
    return out

def split_range(src, start, end, out_dir, prefix="range"):
    os.makedirs(out_dir, exist_ok=True)
    reader = PdfReader(src); total = len(reader.pages)
    start, end = max(1,start), min(total,end)
    w = PdfWriter()
    for i in range(start-1, end): w.add_page(reader.pages[i])
    p = os.path.join(out_dir, f"{prefix}_pages{start}-{end}.pdf")
    with open(p,"wb") as f: w.write(f)
    return p

def split_custom(src, groups, out_dir, prefix="group"):
    os.makedirs(out_dir, exist_ok=True)
    reader = PdfReader(src); total = len(reader.pages); out = []
    for idx, group in enumerate(groups):
        w = PdfWriter()
        valid = [p for p in group if 1 <= p <= total]
        for pg in valid: w.add_page(reader.pages[pg-1])
        if valid:
            label = "_".join(str(p) for p in valid)
            p = os.path.join(out_dir, f"{prefix}{idx+1}_p{label}.pdf")
            with open(p,"wb") as f: w.write(f)
            out.append(p)
    return out

def merge_pdfs(src_list, dst):
    if not HAS_PYPDF: raise ImportError("pip install pypdf")
    w = PdfWriter()
    for src in src_list:
        c = cat(src)
        if c == "pdf":
            for page in PdfReader(src).pages: w.add_page(page)
        elif c == "image":
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tf: tmp = tf.name
            images_to_pdf([src], tmp)
            for page in PdfReader(tmp).pages: w.add_page(page)
            os.unlink(tmp)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst,"wb") as f: w.write(f)
    return dst

def resequence_pdf(src, order, dst):
    reader = PdfReader(src); total = len(reader.pages)
    w = PdfWriter()
    for pg in order:
        if 1 <= pg <= total: w.add_page(reader.pages[pg-1])
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst,"wb") as f: w.write(f)
    return dst

def delete_pages(src, pages_to_delete, dst):
    reader = PdfReader(src); total = len(reader.pages)
    remove = set(pages_to_delete); w = PdfWriter()
    for i in range(total):
        if (i+1) not in remove: w.add_page(reader.pages[i])
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst,"wb") as f: w.write(f)
    return dst

def rotate_pages(src, degrees, pages_to_rotate, dst):
    reader = PdfReader(src); total = len(reader.pages)
    target = set(pages_to_rotate) if pages_to_rotate else set(range(1, total+1))
    w = PdfWriter()
    for i, page in enumerate(reader.pages):
        if (i+1) in target: page.rotate(degrees)
        w.add_page(page)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst,"wb") as f: w.write(f)
    return dst

def reverse_pdf(src, dst):
    reader = PdfReader(src)
    return resequence_pdf(src, list(range(len(reader.pages), 0, -1)), dst)

def compress_pdf(src, dst):
    reader = PdfReader(src); w = PdfWriter()
    for page in reader.pages:
        page.compress_content_streams(); w.add_page(page)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst,"wb") as f: w.write(f)
    return dst

def encrypt_pdf(src, dst, user_pw, owner_pw=""):
    reader = PdfReader(src); w = PdfWriter()
    for page in reader.pages: w.add_page(page)
    w.encrypt(user_pw, owner_pw or user_pw)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst,"wb") as f: w.write(f)
    return dst

def decrypt_pdf(src, dst, password):
    reader = PdfReader(src)
    if reader.is_encrypted:
        if reader.decrypt(password) == 0: raise ValueError("Wrong password")
    w = PdfWriter()
    for page in reader.pages: w.add_page(page)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst,"wb") as f: w.write(f)
    return dst

def watermark_text(src, dst, text, opacity=0.3, fontsize=48,
                   color="#888888", angle=45, pages=None):
    if not (HAS_PIL and HAS_PDF2IMAGE):
        raise ImportError("pip install Pillow pdf2image  (+poppler)")
    reader = PdfReader(src); total = len(reader.pages)
    target = set(pages) if pages else set(range(1, total+1))
    w = PdfWriter(); imgs = convert_from_path(src, dpi=150)
    for i, img in enumerate(imgs):
        pg = i + 1
        if pg in target:
            try:
                font = ImageFont.truetype(
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", fontsize)
            except Exception:
                try: font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", fontsize)
                except: font = ImageFont.load_default()
            try: r,g,b = int(color[1:3],16), int(color[3:5],16), int(color[5:7],16)
            except: r,g,b = 128,128,128
            alpha = int(opacity * 255)
            txt_img = Image.new("RGBA", img.size, (0,0,0,0))
            draw = ImageDraw.Draw(txt_img)
            bbox = draw.textbbox((0,0), text, font=font)
            tw, th = bbox[2]-bbox[0], bbox[3]-bbox[1]
            draw.text((img.size[0]//2-tw//2, img.size[1]//2-th//2),
                      text, font=font, fill=(r,g,b,alpha))
            txt_img = txt_img.rotate(angle, expand=False)
            img = Image.alpha_composite(img.convert("RGBA"), txt_img)
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tf: tmp = tf.name
        img.convert("RGB").save(tmp, "PDF")
        for page in PdfReader(tmp).pages: w.add_page(page)
        os.unlink(tmp)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst,"wb") as f: w.write(f)
    return dst

def watermark_pdf_overlay(src, wm_pdf, dst, pages=None):
    reader = PdfReader(src); wm_page = PdfReader(wm_pdf).pages[0]
    total = len(reader.pages)
    target = set(pages) if pages else set(range(1, total+1))
    w = PdfWriter()
    for i, page in enumerate(reader.pages):
        if (i+1) in target: page.merge_page(wm_page)
        w.add_page(page)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst,"wb") as f: w.write(f)
    return dst

def get_metadata(src):
    reader = PdfReader(src); meta = reader.metadata or {}
    size = os.path.getsize(src)
    return {"File": Path(src).name, "Pages": len(reader.pages),
            "Size": f"{size/1024:.1f} KB  ({size/(1024*1024):.2f} MB)",
            "Title": meta.get("/Title","—"), "Author": meta.get("/Author","—"),
            "Subject": meta.get("/Subject","—"), "Creator": meta.get("/Creator","—"),
            "Producer": meta.get("/Producer","—"),
            "Created": meta.get("/CreationDate","—"),
            "Modified": meta.get("/ModDate","—")}

def set_metadata(src, dst, fields):
    reader = PdfReader(src); w = PdfWriter()
    for page in reader.pages: w.add_page(page)
    existing = dict(reader.metadata or {})
    mapping = {"title":"/Title","author":"/Author","subject":"/Subject","creator":"/Creator"}
    for k, v in fields.items():
        pdf_key = mapping.get(k.lower())
        if pdf_key and v.strip(): existing[pdf_key] = v.strip()
    w.add_metadata(existing)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst,"wb") as f: w.write(f)
    return dst

# ═══════════════════════════════════════════════════════════════════════════════
# ── Master dispatch ───────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

def do_convert(src, out_fmt, out_dir, pages_str="", dpi=150, log=None):
    def L(m): log(m) if log else None
    c = cat(src); stem = Path(src).stem
    out_fmt = out_fmt.lower().lstrip(".")
    os.makedirs(out_dir, exist_ok=True)
    total = pdf_page_count(src) if c == "pdf" and HAS_PYPDF else 0
    pages = parse_pages(pages_str, total) if pages_str.strip() and total else None

    # ── PDF ──────────────────────────────────────────────────────────────────
    if c == "pdf":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "txt":  L("PDF→TXT");   return [pdf_to_txt(src,dst,pages)]
        if out_fmt == "docx": L("PDF→DOCX");  return [pdf_to_docx(src,dst,pages)]
        if out_fmt in ("png","jpg","jpeg"):
            L(f"PDF→{out_fmt.upper()}"); return pdf_to_images(src,os.path.join(out_dir,stem+"_images"),out_fmt,pages,dpi)
        if out_fmt == "html": L("PDF→HTML");  return [pdf_to_html(src,dst,pages)]
        if out_fmt in ("xlsx","xls"): L("PDF→XLSX"); return [pdf_to_xlsx(src,dst.replace(f".{out_fmt}",".xlsx"))]
        if out_fmt == "pptx": L("PDF→PPTX"); return [pdf_to_pptx(src,dst,dpi)]
        if out_fmt == "csv":
            xls = dst.replace(".csv",".xlsx")
            L("PDF→CSV (via xlsx)"); pdf_to_xlsx(src,xls)
            return [excel_to_csv(xls, dst)]

    # ── Excel ────────────────────────────────────────────────────────────────
    if c == "excel":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "pdf":  L("XLSX→PDF");  return [excel_to_pdf(src, dst)]
        if out_fmt == "csv":  L("XLSX→CSV");  return [excel_to_csv(src, dst)]
        if out_fmt == "txt":  L("XLSX→TXT");  return [excel_to_txt(src, dst)]
        if out_fmt == "html": L("XLSX→HTML"); return [excel_to_html(src, dst)]
        if out_fmt == "docx": L("XLSX→DOCX"); return [excel_to_docx(src, dst)]
        if out_fmt in ("xlsx","xls","ods"):   # format-to-format via LibreOffice
            L(f"Excel→{out_fmt.upper()}"); return [lo_convert(src, out_fmt, out_dir)]

    # ── CSV / TSV ─────────────────────────────────────────────────────────────
    if c == "csv":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt in ("xlsx","xls"): L("CSV→XLSX"); return [csv_to_xlsx(src, dst.replace(f".{out_fmt}",".xlsx"))]
        if out_fmt == "pdf":  L("CSV→PDF");  return [csv_to_pdf(src, dst)]
        if out_fmt == "html": L("CSV→HTML"); return [csv_to_html(src, dst)]
        if out_fmt == "txt":  L("CSV→TXT");  return [csv_to_txt(src, dst)]
        if out_fmt == "docx": L("CSV→DOCX"); return [csv_to_docx(src, dst)]

    # ── PPTX ─────────────────────────────────────────────────────────────────
    if c == "pptx":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "pdf":  L("PPTX→PDF");  return [pptx_to_pdf(src, dst)]
        if out_fmt in ("png","jpg","jpeg"):
            L(f"PPTX→{out_fmt.upper()}"); return pptx_to_images(src,os.path.join(out_dir,stem+"_slides"),out_fmt,dpi)
        if out_fmt == "txt":  L("PPTX→TXT");  return [pptx_to_txt(src, dst)]
        if out_fmt == "html": L("PPTX→HTML"); return [pptx_to_html(src, dst)]
        if out_fmt == "docx": L("PPTX→DOCX"); return [pptx_to_docx(src, dst)]

    # ── DOCX ─────────────────────────────────────────────────────────────────
    if c == "docx":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "pdf":  L("DOCX→PDF");  return [docx_to_pdf(src, dst)]
        if out_fmt == "txt":  L("DOCX→TXT");  return [docx_to_txt(src, dst)]
        if out_fmt == "html": L("DOCX→HTML"); return [docx_to_html(src, dst)]

    # ── TXT ──────────────────────────────────────────────────────────────────
    if c == "txt":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "pdf":  L("TXT→PDF");  return [txt_to_pdf(src, dst)]
        if out_fmt == "docx": L("TXT→DOCX"); return [txt_to_docx(src, dst)]
        if out_fmt == "html": L("TXT→HTML"); return [txt_to_html(src, dst)]

    # ── HTML ─────────────────────────────────────────────────────────────────
    if c == "html":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "pdf": L("HTML→PDF"); return [html_to_pdf(src, dst)]
        if out_fmt == "txt": L("HTML→TXT"); return [html_to_txt(src, dst)]

    # ── Image ────────────────────────────────────────────────────────────────
    if c == "image":
        if out_fmt == "pdf":
            dst = os.path.join(out_dir, stem+".pdf")
            L("Image→PDF"); return [images_to_pdf([src],dst)]
        dst = os.path.join(out_dir, stem+"."+out_fmt)
        L(f"Image→{out_fmt.upper()}"); return [image_convert(src,dst)]

    # ── Video / Audio ─────────────────────────────────────────────────────────
    if c in ("video","audio"):
        dst = os.path.join(out_dir, stem+"."+out_fmt)
        L(f"{c.upper()}→{out_fmt.upper()}")
        if out_fmt == "gif":
            return [ffmpeg_convert(src,dst,["-vf","fps=10,scale=480:-1:flags=lanczos","-loop","0"])]
        return [ffmpeg_convert(src,dst)]

    raise ValueError(f"No conversion available: {c} → {out_fmt}")

# ═══════════════════════════════════════════════════════════════════════════════
# COLOURS & FONTS
# ═══════════════════════════════════════════════════════════════════════════════

C = dict(
    bg="#16161e", panel="#1e1e2a", card="#252535", border="#333350",
    accent="#7c6af7", accent2="#f7706a", green="#4ecdc4", yellow="#ffd166",
    red="#f7706a", text="#e8e8f0", dim="#8888aa", dim2="#44445a",
    log_bg="#0e0e18", sidebar="#13131c", sel="#2a2a4a",
)
F = dict(
    title=("Consolas",19,"bold"), head=("Consolas",11,"bold"),
    label=("Consolas",10), body=("Consolas",10), small=("Consolas",9),
    log=("Consolas",9), btn=("Consolas",10,"bold"), btnbig=("Consolas",11,"bold"),
)

class Tooltip:
    def __init__(self, w, text):
        self.w=w; self.text=text; self.tip=None
        w.bind("<Enter>",self.show); w.bind("<Leave>",self.hide)
    def show(self,_=None):
        x=self.w.winfo_rootx()+24; y=self.w.winfo_rooty()+20
        self.tip=tw=tk.Toplevel(self.w); tw.wm_overrideredirect(True)
        tw.wm_geometry(f"+{x}+{y}")
        tk.Label(tw,text=self.text,font=F["small"],bg="#2a2a3e",fg=C["text"],
                 relief="flat",padx=8,pady=4,wraplength=320,justify="left").pack()
    def hide(self,_=None):
        if self.tip: self.tip.destroy(); self.tip=None

# ═══════════════════════════════════════════════════════════════════════════════
# APPLICATION
# ═══════════════════════════════════════════════════════════════════════════════

class App(tk.Tk):
    TABS = [
        ("🔄","Convert","convert"), ("✂️","Split","split"),
        ("🔗","Merge","merge"),     ("📐","Organise","organise"),
        ("💧","Stamp","stamp"),     ("🔒","Protect","protect"),
        ("📦","Compress","compress"),("🏷","Metadata","metadata"),
        ("📋","Queue","queue"),     ("❓","Help","help"),
    ]

    def __init__(self):
        super().__init__()
        self.title("File Workshop v4")
        self.configure(bg=C["bg"])
        self.geometry("1100x760")
        self.minsize(920,640)
        self.files: List[str] = []
        self.out_dir = tk.StringVar()
        self.page_count = 0
        self.active_pdf = ""
        self.pages: Dict[str,tk.Frame] = {}
        self.tab_btns: Dict[str,tk.Button] = {}
        self._build()
        self._switch_tab("convert")
        self._status("Ready — add files to get started")

    # ─── Skeleton ─────────────────────────────────────────────────────────────

    def _build(self):
        hdr = tk.Frame(self, bg=C["sidebar"])
        hdr.pack(fill="x")
        tk.Label(hdr,text="⬛ FILE WORKSHOP",font=F["title"],
                 bg=C["sidebar"],fg=C["accent"]).pack(side="left",padx=20,pady=12)
        tk.Label(hdr,text="convert · split · merge · organise · stamp · protect · excel · pptx",
                 font=("Consolas",9),bg=C["sidebar"],fg=C["dim"]).pack(side="left")
        self.dep_lbl=tk.Label(hdr,text="",font=F["small"],bg=C["sidebar"],fg=C["dim"])
        self.dep_lbl.pack(side="right",padx=18)
        self._refresh_deps()
        body=tk.Frame(self,bg=C["bg"]); body.pack(fill="both",expand=True)
        self._build_sidebar(body); self._build_main(body)

    def _build_sidebar(self, parent):
        sb=tk.Frame(parent,bg=C["sidebar"],width=192)
        sb.pack(side="left",fill="y"); sb.pack_propagate(False)
        tk.Label(sb,text="TOOLS",font=("Consolas",8,"bold"),
                 bg=C["sidebar"],fg=C["dim2"]).pack(anchor="w",padx=14,pady=(16,4))
        for icon,label,key in self.TABS:
            b=tk.Button(sb,text=f" {icon}  {label}",font=F["btn"],
                        bg=C["sidebar"],fg=C["dim"],relief="flat",bd=0,
                        cursor="hand2",anchor="w",padx=14,pady=8,
                        activebackground=C["card"],
                        command=lambda k=key: self._switch_tab(k))
            b.pack(fill="x"); self.tab_btns[key]=b
        tk.Frame(sb,bg=C["border"],height=1).pack(fill="x",padx=10,pady=10)
        tk.Label(sb,text="OUTPUT",font=("Consolas",8,"bold"),
                 bg=C["sidebar"],fg=C["dim2"]).pack(anchor="w",padx=14,pady=(0,4))
        od=tk.Frame(sb,bg=C["sidebar"]); od.pack(fill="x",padx=8,pady=(0,4))
        tk.Entry(od,textvariable=self.out_dir,font=("Consolas",8),
                 bg=C["card"],fg=C["text"],insertbackground=C["text"],
                 relief="flat",bd=0).pack(fill="x",ipady=3,padx=2)
        self._btn(sb,"📁  Set Folder",self._browse_output,C["dim2"]).pack(fill="x",padx=8,pady=2)
        self._btn(sb,"📂  Open Output",self._open_output,C["dim2"]).pack(fill="x",padx=8,pady=2)
        tk.Frame(sb,bg=C["border"],height=1).pack(fill="x",padx=10,pady=8)
        tk.Label(sb,text="FILES",font=("Consolas",8,"bold"),
                 bg=C["sidebar"],fg=C["dim2"]).pack(anchor="w",padx=14,pady=(0,4))
        self._btn(sb,"+ Add Files",self._add_files,C["accent"]).pack(fill="x",padx=8,pady=2)
        self._btn(sb,"+ Add Folder",self._add_folder,C["dim2"]).pack(fill="x",padx=8,pady=2)
        self._btn(sb,"🗑  Clear All",self._clear_queue,C["dim2"]).pack(fill="x",padx=8,pady=2)
        self.pdf_badge=tk.Label(sb,text="",font=("Consolas",8),
                                bg=C["sidebar"],fg=C["green"],
                                wraplength=170,justify="left")
        self.pdf_badge.pack(anchor="w",padx=10,pady=(8,0))

    def _build_main(self, parent):
        main=tk.Frame(parent,bg=C["bg"]); main.pack(side="left",fill="both",expand=True)
        for _,_,key in self.TABS:
            self.pages[key]=tk.Frame(main,bg=C["bg"])
        self._build_convert(); self._build_split(); self._build_merge()
        self._build_organise(); self._build_stamp(); self._build_protect()
        self._build_compress(); self._build_metadata()
        self._build_queue(); self._build_help()
        # Log
        log_wrap=tk.Frame(main,bg=C["log_bg"])
        log_wrap.pack(fill="x",padx=14,pady=(0,8),side="bottom")
        lhdr=tk.Frame(log_wrap,bg=C["log_bg"]); lhdr.pack(fill="x")
        tk.Label(lhdr,text="LOG",font=("Consolas",8,"bold"),
                 bg=C["log_bg"],fg=C["dim2"]).pack(side="left",padx=8,pady=(4,2))
        self._btn(lhdr,"Clear",self._clear_log,C["dim2"]).pack(side="right",padx=6)
        self.log_box=scrolledtext.ScrolledText(
            log_wrap,height=6,font=F["log"],bg=C["log_bg"],fg=C["text"],
            insertbackground=C["text"],relief="flat",state="disabled",bd=0)
        self.log_box.pack(fill="x",padx=2,pady=(0,4))
        for tag,col in [("ok",C["green"]),("err",C["red"]),("info",C["accent"]),("warn",C["yellow"])]:
            self.log_box.tag_config(tag,foreground=col)
        self.status_var=tk.StringVar(value="Ready")
        tk.Label(main,textvariable=self.status_var,font=F["small"],
                 bg=C["border"],fg=C["dim"],anchor="w",padx=10,pady=3).pack(fill="x",side="bottom")

    # ─── Convert tab ──────────────────────────────────────────────────────────

    def _build_convert(self):
        f=self.pages["convert"]
        self._section_title(f,"CONVERT FILES","Batch-convert any supported file to another format")
        fmt_card=self._card(f)
        self._lbl(fmt_card,"OUTPUT FORMAT",head=True).pack(anchor="w",pady=(0,10))
        cats=[
            ("Documents", ["txt","docx","pdf","html"]),
            ("Spreadsheet",["xlsx","csv"]),
            ("Presentation",["pptx"]),
            ("Images",    ["png","jpg","webp","bmp","gif","tiff","ico"]),
            ("Audio",     ["mp3","wav","ogg","flac","aac","m4a"]),
            ("Video",     ["mp4","avi","mov","mkv","webm","gif"]),
        ]
        self.out_fmt=tk.StringVar(value="pdf")
        for cat_name,fmts in cats:
            row=tk.Frame(fmt_card,bg=C["panel"]); row.pack(fill="x",pady=2)
            self._lbl(row,f"{cat_name}:",dim=True,w=13).pack(side="left")
            for fmt in fmts:
                rb=tk.Radiobutton(row,text=fmt.upper(),variable=self.out_fmt,
                                  value=fmt,font=F["btn"],bg=C["card"],fg=C["dim"],
                                  selectcolor=C["accent"],activebackground=C["card"],
                                  indicatoron=0,relief="flat",padx=8,pady=3,
                                  cursor="hand2",command=self._on_fmt_change)
                rb.pack(side="left",padx=3)
                self.out_fmt.trace_add("write",
                    lambda *_,rb=rb,fmt=fmt:
                        rb.config(bg=C["accent"],fg="#fff")
                        if self.out_fmt.get()==fmt
                        else rb.config(bg=C["card"],fg=C["dim"]))
        opt_card=self._card(f)
        self._lbl(opt_card,"OPTIONS",head=True).pack(anchor="w",pady=(0,8))
        r1=tk.Frame(opt_card,bg=C["panel"]); r1.pack(fill="x",pady=2)
        self._lbl(r1,"Pages (PDF/PPTX):",w=18).pack(side="left")
        self.conv_pages=tk.StringVar()
        e=tk.Entry(r1,textvariable=self.conv_pages,font=F["body"],
                   bg=C["card"],fg=C["text"],insertbackground=C["text"],
                   relief="flat",bd=0,width=22)
        e.pack(side="left",ipady=3,padx=(6,6))
        self._lbl(r1,'e.g. "1,3,5-8" — blank = all',dim=True).pack(side="left")
        r2=tk.Frame(opt_card,bg=C["panel"]); r2.pack(fill="x",pady=6)
        self._lbl(r2,"DPI (image export):",w=18).pack(side="left")
        self.conv_dpi=tk.StringVar(value="150")
        for val,lbl in [("72","72 draft"),("150","150 normal"),("300","300 print")]:
            tk.Radiobutton(r2,text=lbl,variable=self.conv_dpi,value=val,
                           font=F["small"],bg=C["panel"],fg=C["dim"],
                           selectcolor=C["accent"],activebackground=C["panel"],
                           cursor="hand2").pack(side="left",padx=(0,14))
        self.conv_note=tk.Label(opt_card,text="",font=F["small"],
                                bg=C["panel"],fg=C["yellow"],wraplength=680,justify="left")
        self.conv_note.pack(anchor="w",pady=(4,0))
        self._bigbtn(f,"▶   RUN CONVERSION",self._run_convert).pack(pady=12)

    # ─── Split tab ────────────────────────────────────────────────────────────

    def _build_split(self):
        f=self.pages["split"]
        self._section_title(f,"SPLIT","Break a PDF or PPTX into separate files")
        type_card=self._card(f)
        self._lbl(type_card,"FILE TYPE",head=True).pack(anchor="w",pady=(0,6))
        self.split_type=tk.StringVar(value="pdf")
        for val,lbl in [("pdf","PDF — split pages"),("pptx","PPTX — split slides")]:
            tk.Radiobutton(type_card,text=lbl,variable=self.split_type,value=val,
                           font=F["body"],bg=C["panel"],fg=C["text"],
                           selectcolor=C["accent2"],activebackground=C["panel"],
                           command=self._update_split_ui,cursor="hand2").pack(anchor="w",pady=2)
        self.split_info_lbl=tk.Label(type_card,text="No PDF/PPTX loaded — add one to the queue",
                                      font=F["small"],bg=C["panel"],fg=C["dim"])
        self.split_info_lbl.pack(anchor="w",pady=(4,0))
        mode_card=self._card(f)
        self._lbl(mode_card,"MODE",head=True).pack(anchor="w",pady=(0,8))
        self.split_mode=tk.StringVar(value="each")
        for val,lbl,tip in [
            ("each",  "Each page/slide → own file", "10-page PDF becomes 10 files"),
            ("range", "Range → one file",            "Pages 3–8 into one file"),
            ("custom","Custom groups (PDF only)",    "1,10 | 3,5 | 2-4,7"),
        ]:
            r=tk.Frame(mode_card,bg=C["panel"]); r.pack(anchor="w",pady=2)
            tk.Radiobutton(r,text=lbl,variable=self.split_mode,value=val,
                           font=F["body"],bg=C["panel"],fg=C["text"],
                           selectcolor=C["accent2"],activebackground=C["panel"],
                           command=self._update_split_ui,cursor="hand2").pack(side="left")
            self._lbl(r,f"  — {tip}",dim=True).pack(side="left")
        self.split_opts=self._card(f); self._update_split_ui()
        pr=tk.Frame(f,bg=C["bg"]); pr.pack(fill="x",padx=16,pady=2)
        self._lbl(pr,"File prefix:",w=12).pack(side="left")
        self.split_prefix=tk.StringVar(value="split")
        tk.Entry(pr,textvariable=self.split_prefix,font=F["body"],
                 bg=C["card"],fg=C["text"],insertbackground=C["text"],
                 relief="flat",bd=0,width=16).pack(side="left",padx=(6,0),ipady=3)
        self._bigbtn(f,"▶   RUN SPLIT",self._run_split,C["accent2"]).pack(pady=12)

    # ─── Merge tab ────────────────────────────────────────────────────────────

    def _build_merge(self):
        f=self.pages["merge"]
        self._section_title(f,"MERGE","Combine multiple files into one")
        type_card=self._card(f)
        self._lbl(type_card,"MERGE TYPE",head=True).pack(anchor="w",pady=(0,6))
        self.merge_type=tk.StringVar(value="pdf")
        for val,lbl,tip in [
            ("pdf", "Merge → PDF",  "Combine PDFs and/or images into one PDF"),
            ("pptx","Merge → PPTX", "Combine multiple PPTX files into one presentation"),
        ]:
            r=tk.Frame(type_card,bg=C["panel"]); r.pack(anchor="w",pady=2)
            tk.Radiobutton(r,text=lbl,variable=self.merge_type,value=val,
                           font=F["body"],bg=C["panel"],fg=C["text"],
                           selectcolor=C["accent"],activebackground=C["panel"],
                           cursor="hand2").pack(side="left")
            self._lbl(r,f"  — {tip}",dim=True).pack(side="left")
        opt_card=self._card(f)
        self._lbl(opt_card,"OUTPUT FILENAME",head=True).pack(anchor="w",pady=(0,8))
        row=tk.Frame(opt_card,bg=C["panel"]); row.pack(fill="x")
        self._lbl(row,"Filename:",w=12).pack(side="left")
        self.merge_name=tk.StringVar(value="merged_output")
        tk.Entry(row,textvariable=self.merge_name,font=F["body"],
                 bg=C["card"],fg=C["text"],insertbackground=C["text"],
                 relief="flat",bd=0,width=30).pack(side="left",padx=(6,6),ipady=3)
        self.merge_ext_lbl=self._lbl(row,".pdf",dim=True)
        self.merge_ext_lbl.pack(side="left")
        self.merge_type.trace_add("write",lambda *_: self.merge_ext_lbl.config(
            text="."+self.merge_type.get()))
        preview_card=self._card(f)
        self._lbl(preview_card,"FILES IN MERGE ORDER",head=True).pack(anchor="w",pady=(0,6))
        self.merge_preview=tk.Text(preview_card,height=6,font=F["small"],
                                    bg=C["log_bg"],fg=C["dim"],relief="flat",
                                    state="disabled",bd=0)
        self.merge_preview.pack(fill="x")
        self._btn(preview_card,"↻ Refresh",self._refresh_merge_preview,C["dim2"]).pack(anchor="w",pady=(6,0))
        self._bigbtn(f,"▶   RUN MERGE",self._run_merge).pack(pady=12)

    # ─── Organise tab ─────────────────────────────────────────────────────────

    def _build_organise(self):
        f=self.pages["organise"]
        self._section_title(f,"ORGANISE PDF PAGES","Resequence · Delete · Rotate · Reverse")
        lbl_card=self._card(f)
        self.org_pdf_lbl_var=tk.StringVar(value="No PDF — add one to the queue")
        tk.Label(lbl_card,textvariable=self.org_pdf_lbl_var,font=F["small"],
                 bg=C["panel"],fg=C["green"]).pack(anchor="w")
        op_card=self._card(f)
        self._lbl(op_card,"OPERATION",head=True).pack(anchor="w",pady=(0,8))
        self.org_op=tk.StringVar(value="resequence")
        for val,lbl,tip in [
            ("resequence","Resequence","Custom order, e.g. 3,1,2"),
            ("delete",    "Delete pages","Remove pages, e.g. 2,5,7-9"),
            ("rotate",    "Rotate pages","90/180/270° on all or specific pages"),
            ("reverse",   "Reverse","Flip entire page order"),
        ]:
            r=tk.Frame(op_card,bg=C["panel"]); r.pack(anchor="w",pady=2)
            tk.Radiobutton(r,text=lbl,variable=self.org_op,value=val,
                           font=F["body"],bg=C["panel"],fg=C["text"],
                           selectcolor=C["accent"],activebackground=C["panel"],
                           command=self._update_org_ui,cursor="hand2").pack(side="left")
            self._lbl(r,f"  — {tip}",dim=True).pack(side="left")
        self.org_opts=self._card(f); self._update_org_ui()
        row=tk.Frame(f,bg=C["bg"]); row.pack(fill="x",padx=16,pady=2)
        self._lbl(row,"Output name:",w=14).pack(side="left")
        self.org_out_name=tk.StringVar(value="organised")
        tk.Entry(row,textvariable=self.org_out_name,font=F["body"],
                 bg=C["card"],fg=C["text"],insertbackground=C["text"],
                 relief="flat",bd=0,width=24).pack(side="left",padx=(6,6),ipady=3)
        self._lbl(row,".pdf",dim=True).pack(side="left")
        self._bigbtn(f,"▶   RUN ORGANISE",self._run_organise,C["accent"]).pack(pady=12)

    # ─── Stamp tab ────────────────────────────────────────────────────────────

    def _build_stamp(self):
        f=self.pages["stamp"]
        self._section_title(f,"STAMP / WATERMARK","Overlay text or PDF watermark on pages")
        mode_card=self._card(f)
        self._lbl(mode_card,"STAMP TYPE",head=True).pack(anchor="w",pady=(0,6))
        self.stamp_mode=tk.StringVar(value="text")
        for val,lbl in [("text","Text watermark"),("pdf","PDF overlay")]:
            tk.Radiobutton(mode_card,text=lbl,variable=self.stamp_mode,value=val,
                           font=F["body"],bg=C["panel"],fg=C["text"],
                           selectcolor=C["accent"],activebackground=C["panel"],
                           command=self._update_stamp_ui,cursor="hand2").pack(anchor="w",pady=2)
        self.stamp_opts=self._card(f); self._update_stamp_ui()
        pg_card=self._card(f)
        self._lbl(pg_card,"APPLY TO PAGES",head=True).pack(anchor="w",pady=(0,6))
        row=tk.Frame(pg_card,bg=C["panel"]); row.pack(fill="x")
        self._lbl(row,"Pages:",w=8).pack(side="left")
        self.stamp_pages=tk.StringVar()
        tk.Entry(row,textvariable=self.stamp_pages,font=F["body"],
                 bg=C["card"],fg=C["text"],insertbackground=C["text"],
                 relief="flat",bd=0,width=22).pack(side="left",padx=(6,6),ipady=3)
        self._lbl(row,"blank = all",dim=True).pack(side="left")
        r2=tk.Frame(f,bg=C["bg"]); r2.pack(fill="x",padx=16,pady=2)
        self._lbl(r2,"Output name:",w=14).pack(side="left")
        self.stamp_out_name=tk.StringVar(value="stamped")
        tk.Entry(r2,textvariable=self.stamp_out_name,font=F["body"],
                 bg=C["card"],fg=C["text"],insertbackground=C["text"],
                 relief="flat",bd=0,width=24).pack(side="left",padx=(6,6),ipady=3)
        self._lbl(r2,".pdf",dim=True).pack(side="left")
        self._bigbtn(f,"▶   RUN STAMP",self._run_stamp,C["accent"]).pack(pady=12)

    # ─── Protect tab ──────────────────────────────────────────────────────────

    def _build_protect(self):
        f=self.pages["protect"]
        self._section_title(f,"PROTECT / DECRYPT PDF","Password-protect or remove password")
        mode_card=self._card(f)
        self.protect_mode=tk.StringVar(value="encrypt")
        for val,lbl,tip in [
            ("encrypt","Encrypt (add password)","Lock with password"),
            ("decrypt","Decrypt (remove password)","Requires current password"),
        ]:
            r=tk.Frame(mode_card,bg=C["panel"]); r.pack(anchor="w",pady=2)
            tk.Radiobutton(r,text=lbl,variable=self.protect_mode,value=val,
                           font=F["body"],bg=C["panel"],fg=C["text"],
                           selectcolor=C["accent"],activebackground=C["panel"],
                           cursor="hand2").pack(side="left")
            self._lbl(r,f"  — {tip}",dim=True).pack(side="left")
        pw_card=self._card(f)
        self._lbl(pw_card,"PASSWORD",head=True).pack(anchor="w",pady=(0,8))
        for lbl,attr in [("User password:","protect_pw1"),("Owner password (optional):","protect_pw2")]:
            r=tk.Frame(pw_card,bg=C["panel"]); r.pack(fill="x",pady=3)
            self._lbl(r,lbl,w=24).pack(side="left")
            var=tk.StringVar(); setattr(self,attr,var)
            tk.Entry(r,textvariable=var,font=F["body"],bg=C["card"],fg=C["text"],
                     insertbackground=C["text"],show="•",relief="flat",bd=0,width=26
                     ).pack(side="left",padx=(6,0),ipady=3)
        row=tk.Frame(f,bg=C["bg"]); row.pack(fill="x",padx=16,pady=2)
        self._lbl(row,"Output name:",w=14).pack(side="left")
        self.protect_out=tk.StringVar(value="protected")
        tk.Entry(row,textvariable=self.protect_out,font=F["body"],
                 bg=C["card"],fg=C["text"],insertbackground=C["text"],
                 relief="flat",bd=0,width=24).pack(side="left",padx=(6,6),ipady=3)
        self._lbl(row,".pdf",dim=True).pack(side="left")
        self._bigbtn(f,"▶   RUN",self._run_protect,C["accent"]).pack(pady=12)

    # ─── Compress tab ─────────────────────────────────────────────────────────

    def _build_compress(self):
        f=self.pages["compress"]
        self._section_title(f,"COMPRESS PDF","Reduce file size via lossless stream compression")
        info=self._card(f)
        self._lbl(info,"Compresses content streams. Best for text-heavy PDFs. "
                  "For image-heavy PDFs, re-export at lower DPI from the Convert tab.",
                  dim=True,wrap=700).pack(anchor="w")
        row=tk.Frame(f,bg=C["bg"]); row.pack(fill="x",padx=16,pady=8)
        self._lbl(row,"Output name:",w=14).pack(side="left")
        self.compress_out=tk.StringVar(value="compressed")
        tk.Entry(row,textvariable=self.compress_out,font=F["body"],
                 bg=C["card"],fg=C["text"],insertbackground=C["text"],
                 relief="flat",bd=0,width=24).pack(side="left",padx=(6,6),ipady=3)
        self._lbl(row,".pdf",dim=True).pack(side="left")
        self._bigbtn(f,"▶   RUN COMPRESS",self._run_compress,C["accent"]).pack(pady=12)

    # ─── Metadata tab ─────────────────────────────────────────────────────────

    def _build_metadata(self):
        f=self.pages["metadata"]
        self._section_title(f,"PDF METADATA","View and edit document properties")
        view_card=self._card(f)
        self._lbl(view_card,"CURRENT METADATA",head=True).pack(anchor="w",pady=(0,6))
        self.meta_display=tk.Text(view_card,height=9,font=F["body"],
                                   bg=C["log_bg"],fg=C["text"],
                                   relief="flat",state="disabled",bd=0)
        self.meta_display.pack(fill="x")
        self._btn(view_card,"↻ Load from active PDF",self._load_metadata,C["dim2"]).pack(anchor="w",pady=(6,0))
        edit_card=self._card(f)
        self._lbl(edit_card,"EDIT FIELDS",head=True).pack(anchor="w",pady=(0,8))
        self.meta_fields: Dict[str,tk.StringVar]={}
        for field in ["Title","Author","Subject","Creator"]:
            r=tk.Frame(edit_card,bg=C["panel"]); r.pack(fill="x",pady=3)
            self._lbl(r,field+":",w=10).pack(side="left")
            var=tk.StringVar(); self.meta_fields[field.lower()]=var
            tk.Entry(r,textvariable=var,font=F["body"],bg=C["card"],fg=C["text"],
                     insertbackground=C["text"],relief="flat",bd=0,width=44
                     ).pack(side="left",padx=(6,0),ipady=3)
        row=tk.Frame(f,bg=C["bg"]); row.pack(fill="x",padx=16,pady=2)
        self._lbl(row,"Output name:",w=14).pack(side="left")
        self.meta_out=tk.StringVar(value="updated_metadata")
        tk.Entry(row,textvariable=self.meta_out,font=F["body"],
                 bg=C["card"],fg=C["text"],insertbackground=C["text"],
                 relief="flat",bd=0,width=24).pack(side="left",padx=(6,6),ipady=3)
        self._lbl(row,".pdf",dim=True).pack(side="left")
        self._bigbtn(f,"▶   SAVE METADATA",self._run_metadata,C["accent"]).pack(pady=12)

    # ─── Queue tab ────────────────────────────────────────────────────────────

    def _build_queue(self):
        f=self.pages["queue"]
        self._section_title(f,"FILE QUEUE","All staged files — reorder with ↑/↓, Ctrl+click for multi-select")
        lf=tk.Frame(f,bg=C["panel"]); lf.pack(fill="both",expand=True,padx=14,pady=(0,6))
        self.queue_lb=tk.Listbox(lf,font=F["body"],bg=C["card"],fg=C["text"],
                                  selectbackground=C["accent"],activestyle="none",
                                  relief="flat",bd=0,selectmode="extended")
        sb=tk.Scrollbar(lf,orient="vertical",command=self.queue_lb.yview)
        self.queue_lb.config(yscrollcommand=sb.set)
        self.queue_lb.pack(side="left",fill="both",expand=True)
        sb.pack(side="right",fill="y")
        br=tk.Frame(f,bg=C["bg"]); br.pack(fill="x",padx=14,pady=(0,4))
        for lbl,cmd in [("↑ Up",self._move_up),("↓ Down",self._move_down),
                          ("✖ Remove",self._remove_selected),("🗑 Clear",self._clear_queue)]:
            self._btn(br,lbl,cmd,C["dim2"]).pack(side="left",padx=(0,6))

    # ─── Help tab ─────────────────────────────────────────────────────────────

    def _build_help(self):
        f=self.pages["help"]
        txt=scrolledtext.ScrolledText(f,font=("Consolas",10),bg=C["panel"],fg=C["text"],
                                       relief="flat",state="normal",padx=18,pady=14,wrap="word")
        txt.pack(fill="both",expand=True,padx=14,pady=10)
        txt.insert("1.0",HELP_TEXT); txt.config(state="disabled")

    # ─── Dynamic sub-UIs ──────────────────────────────────────────────────────

    def _update_split_ui(self):
        for w in self.split_opts.winfo_children(): w.destroy()
        mode=self.split_mode.get(); f=self.split_opts
        if mode=="each":
            self._lbl(f,"Every page/slide becomes its own numbered file.").pack(anchor="w")
        elif mode=="range":
            r=tk.Frame(f,bg=C["panel"]); r.pack(anchor="w")
            self._lbl(r,"Start:",w=7).pack(side="left")
            self.range_start=tk.StringVar(value="1")
            tk.Entry(r,textvariable=self.range_start,font=F["body"],bg=C["card"],
                     fg=C["text"],insertbackground=C["text"],relief="flat",bd=0,width=5
                     ).pack(side="left",padx=(4,12),ipady=3)
            self._lbl(r,"End:",w=5).pack(side="left")
            self.range_end=tk.StringVar(value="5")
            tk.Entry(r,textvariable=self.range_end,font=F["body"],bg=C["card"],
                     fg=C["text"],insertbackground=C["text"],relief="flat",bd=0,width=5
                     ).pack(side="left",padx=(4,0),ipady=3)
        elif mode=="custom":
            self._lbl(f,"Groups by  |  ·  Pages by  ,  ·  Ranges with  -  (PDF only)",dim=True).pack(anchor="w")
            self._lbl(f,"Example:  1,10 | 3,5 | 2-4,7",head=True).pack(anchor="w",pady=(2,4))
            self.custom_groups=tk.StringVar(value="1,2 | 3,4 | 5-7")
            tk.Entry(f,textvariable=self.custom_groups,font=F["body"],bg=C["card"],
                     fg=C["text"],insertbackground=C["text"],relief="flat",bd=0
                     ).pack(fill="x",ipady=4)

    def _update_org_ui(self):
        for w in self.org_opts.winfo_children(): w.destroy()
        op=self.org_op.get(); f=self.org_opts
        if op=="resequence":
            self._lbl(f,"New page order — pages not listed are DROPPED.",dim=True).pack(anchor="w")
            self._lbl(f,"Example:  3,1,2  →  page 3 first, then 1, then 2",head=True).pack(anchor="w",pady=(2,4))
            self.org_seq=tk.StringVar(value="1,2,3")
            tk.Entry(f,textvariable=self.org_seq,font=F["body"],bg=C["card"],
                     fg=C["text"],insertbackground=C["text"],relief="flat",bd=0
                     ).pack(fill="x",ipady=4)
        elif op=="delete":
            self._lbl(f,"Pages to DELETE — all others are kept.",dim=True).pack(anchor="w")
            self._lbl(f,"Example:  2,5,7-9",head=True).pack(anchor="w",pady=(2,4))
            self.org_del_pages=tk.StringVar(value="")
            tk.Entry(f,textvariable=self.org_del_pages,font=F["body"],bg=C["card"],
                     fg=C["text"],insertbackground=C["text"],relief="flat",bd=0
                     ).pack(fill="x",ipady=4)
        elif op=="rotate":
            r=tk.Frame(f,bg=C["panel"]); r.pack(fill="x",pady=2)
            self._lbl(r,"Degrees:",w=10).pack(side="left")
            self.org_rotate_deg=tk.StringVar(value="90")
            for d in ["90","180","270"]:
                tk.Radiobutton(r,text=d+"°",variable=self.org_rotate_deg,value=d,
                               font=F["btn"],bg=C["panel"],fg=C["dim"],
                               selectcolor=C["accent"],activebackground=C["panel"],
                               cursor="hand2").pack(side="left",padx=(0,12))
            r2=tk.Frame(f,bg=C["panel"]); r2.pack(fill="x",pady=(4,0))
            self._lbl(r2,"Pages (blank=all):",w=18).pack(side="left")
            self.org_rotate_pages=tk.StringVar(value="")
            tk.Entry(r2,textvariable=self.org_rotate_pages,font=F["body"],bg=C["card"],
                     fg=C["text"],insertbackground=C["text"],relief="flat",bd=0,width=22
                     ).pack(side="left",padx=(6,0),ipady=3)
        elif op=="reverse":
            self._lbl(f,"Entire page order will be reversed (last page → first).").pack(anchor="w")

    def _update_stamp_ui(self):
        for w in self.stamp_opts.winfo_children(): w.destroy()
        mode=self.stamp_mode.get(); f=self.stamp_opts
        if mode=="text":
            self._lbl(f,"TEXT WATERMARK",head=True).pack(anchor="w",pady=(0,8))
            for lbl,attr,default in [
                ("Watermark text:","stamp_text","CONFIDENTIAL"),
                ("Font size:",     "stamp_size","48"),
                ("Color (hex):",   "stamp_color","#888888"),
                ("Opacity (0–1):","stamp_opacity","0.3"),
                ("Angle (°):",    "stamp_angle","45"),
            ]:
                r=tk.Frame(f,bg=C["panel"]); r.pack(fill="x",pady=2)
                self._lbl(r,lbl,w=18).pack(side="left")
                var=tk.StringVar(value=default); setattr(self,attr,var)
                tk.Entry(r,textvariable=var,font=F["body"],bg=C["card"],fg=C["text"],
                         insertbackground=C["text"],relief="flat",bd=0,width=22
                         ).pack(side="left",padx=(6,0),ipady=3)
        else:
            self._lbl(f,"PDF OVERLAY — first page of selected PDF becomes the stamp",head=True).pack(anchor="w",pady=(0,8))
            r=tk.Frame(f,bg=C["panel"]); r.pack(fill="x")
            self.stamp_wm_path=tk.StringVar(value="")
            tk.Entry(r,textvariable=self.stamp_wm_path,font=F["body"],bg=C["card"],
                     fg=C["text"],insertbackground=C["text"],relief="flat",bd=0
                     ).pack(side="left",fill="x",expand=True,ipady=3,padx=(0,6))
            self._btn(r,"Browse",self._browse_stamp_pdf,C["dim2"]).pack(side="left")

    # ─── Run handlers ─────────────────────────────────────────────────────────

    def _require_pdf(self,label=""):
        pdfs=[f for f in self.files if cat(f)=="pdf"]
        if not pdfs: messagebox.showwarning(label or "No PDF","Add a PDF to the queue first."); return None
        return pdfs[0]

    def _require_out(self):
        out=self.out_dir.get()
        if not out: messagebox.showwarning("No Output","Set an output folder first."); return None
        return out

    def _run_convert(self):
        if not self.files: messagebox.showwarning("No files","Add files first."); return
        out=self._require_out()
        if not out: return
        fmt=self.out_fmt.get()
        pages_str=self.conv_pages.get().strip()
        try: dpi=int(self.conv_dpi.get())
        except: dpi=150
        def task():
            total=len(self.files); ok=0
            for idx,src in enumerate(list(self.files)):
                self._status(f"Converting {idx+1}/{total}: {Path(src).name}")
                try:
                    results=do_convert(src,fmt,out,pages_str,dpi,log=self._log)
                    for r in results: self._log(f"→ {Path(r).name}","ok")
                    ok+=1
                except Exception as e:
                    self._log(f"✖ {Path(src).name}: {e}","err")
            self._log(f"Done — {ok}/{total} converted.","info")
            self._status(f"Done — {ok}/{total} converted")
            self.after(300,self._open_output)
        threading.Thread(target=task,daemon=True).start()

    def _run_split(self):
        out=self._require_out()
        if not out: return
        stype=self.split_type.get()
        mode=self.split_mode.get()
        prefix=self.split_prefix.get().strip() or "split"

        if stype=="pptx":
            pptxs=[f for f in self.files if cat(f)=="pptx"]
            if not pptxs: messagebox.showwarning("No PPTX","Add a PPTX to the queue first."); return
            src=pptxs[0]
            def task():
                try:
                    self._log(f"Splitting slides of {Path(src).name} …")
                    files=pptx_split_slides(src,out,prefix)
                    self._log(f"✔ {len(files)} slide files created.","ok")
                    self.after(300,self._open_output)
                except Exception as ex: self._log(f"ERROR: {ex}","err")
            threading.Thread(target=task,daemon=True).start()
            return

        # PDF split
        src=self._require_pdf("Split")
        if not src: return
        def task():
            try:
                if mode=="each":
                    self._log(f"Splitting every page of {Path(src).name} …")
                    files=split_each(src,out,prefix)
                    self._log(f"✔ {len(files)} files created.","ok")
                elif mode=="range":
                    try: s,e=int(self.range_start.get()),int(self.range_end.get())
                    except: self._log("⚠ Invalid range.","err"); return
                    self._log(f"Extracting pages {s}–{e} …")
                    p=split_range(src,s,e,out,prefix)
                    self._log(f"✔ Saved: {Path(p).name}","ok")
                elif mode=="custom":
                    raw=self.custom_groups.get()
                    groups=parse_groups(raw,self.page_count or 9999)
                    if not groups: self._log("⚠ No valid groups.","err"); return
                    files=split_custom(src,groups,out,prefix)
                    for fp in files: self._log(f"  → {Path(fp).name}")
                    self._log(f"✔ {len(files)} files created.","ok")
                self.after(300,self._open_output)
            except Exception as ex: self._log(f"ERROR: {ex}","err")
        threading.Thread(target=task,daemon=True).start()

    def _run_merge(self):
        if not self.files: messagebox.showwarning("Merge","Add files to the queue first."); return
        out=self._require_out()
        if not out: return
        mtype=self.merge_type.get()
        name=(self.merge_name.get().strip() or "merged_output")
        dst=os.path.join(out,name+"."+mtype)
        srcs=list(self.files)
        def task():
            try:
                self._log(f"Merging {len(srcs)} files into {mtype.upper()} …")
                if mtype=="pdf": merge_pdfs(srcs,dst)
                else: pptx_merge(srcs,dst)
                sz=os.path.getsize(dst)/1024
                self._log(f"✔ {Path(dst).name}  ({sz:.1f} KB)","ok")
                self.after(300,self._open_output)
            except Exception as ex: self._log(f"ERROR: {ex}","err")
        threading.Thread(target=task,daemon=True).start()

    def _run_organise(self):
        src=self._require_pdf("Organise"); out=self._require_out()
        if not src or not out: return
        op=self.org_op.get()
        name=(self.org_out_name.get().strip() or "organised")+".pdf"
        dst=os.path.join(out,name)
        pc=self.page_count or pdf_page_count(src)
        def task():
            try:
                if op=="resequence":
                    order=parse_pages(self.org_seq.get(),pc*10)
                    order=[p for p in order if 1<=p<=pc]
                    if not order: self._log("⚠ No valid order.","err"); return
                    resequence_pdf(src,order,dst)
                elif op=="delete":
                    pages=parse_pages(self.org_del_pages.get(),pc)
                    if not pages: self._log("⚠ No pages specified.","err"); return
                    delete_pages(src,pages,dst)
                elif op=="rotate":
                    try: deg=int(self.org_rotate_deg.get())
                    except: deg=90
                    pstr=self.org_rotate_pages.get().strip()
                    pages=parse_pages(pstr,pc) if pstr else None
                    rotate_pages(src,deg,pages,dst)
                elif op=="reverse":
                    reverse_pdf(src,dst)
                sz=os.path.getsize(dst)/1024
                self._log(f"✔ {Path(dst).name}  ({sz:.1f} KB)","ok")
                self.after(300,self._open_output)
            except Exception as ex: self._log(f"ERROR: {ex}","err")
        threading.Thread(target=task,daemon=True).start()

    def _run_stamp(self):
        src=self._require_pdf("Stamp"); out=self._require_out()
        if not src or not out: return
        mode=self.stamp_mode.get()
        name=(self.stamp_out_name.get().strip() or "stamped")+".pdf"
        dst=os.path.join(out,name)
        pc=self.page_count or pdf_page_count(src)
        pstr=self.stamp_pages.get().strip()
        pages=parse_pages(pstr,pc) if pstr else None
        def task():
            try:
                if mode=="text":
                    text=self.stamp_text.get().strip() or "WATERMARK"
                    try: fsize=int(self.stamp_size.get())
                    except: fsize=48
                    color=self.stamp_color.get().strip() or "#888888"
                    try: opacity=float(self.stamp_opacity.get())
                    except: opacity=0.3
                    try: angle=int(self.stamp_angle.get())
                    except: angle=45
                    self._log(f"Applying text watermark '{text}' …")
                    watermark_text(src,dst,text,opacity,fsize,color,angle,pages)
                else:
                    wm=self.stamp_wm_path.get().strip()
                    if not wm or not os.path.isfile(wm):
                        self._log("⚠ Select a watermark PDF first.","err"); return
                    self._log("Applying PDF overlay …")
                    watermark_pdf_overlay(src,wm,dst,pages)
                sz=os.path.getsize(dst)/1024
                self._log(f"✔ {Path(dst).name}  ({sz:.1f} KB)","ok")
                self.after(300,self._open_output)
            except Exception as ex: self._log(f"ERROR: {ex}","err")
        threading.Thread(target=task,daemon=True).start()

    def _run_protect(self):
        src=self._require_pdf("Protect"); out=self._require_out()
        if not src or not out: return
        mode=self.protect_mode.get()
        pw1=self.protect_pw1.get(); pw2=self.protect_pw2.get()
        dst=os.path.join(out,(self.protect_out.get().strip() or "protected")+".pdf")
        def task():
            try:
                if mode=="encrypt":
                    if not pw1: self._log("⚠ Enter a password.","err"); return
                    encrypt_pdf(src,dst,pw1,pw2); self._log(f"✔ {Path(dst).name}","ok")
                else:
                    if not pw1: self._log("⚠ Enter the password.","err"); return
                    decrypt_pdf(src,dst,pw1); self._log(f"✔ {Path(dst).name}","ok")
                self.after(300,self._open_output)
            except Exception as ex: self._log(f"ERROR: {ex}","err")
        threading.Thread(target=task,daemon=True).start()

    def _run_compress(self):
        src=self._require_pdf("Compress"); out=self._require_out()
        if not src or not out: return
        dst=os.path.join(out,(self.compress_out.get().strip() or "compressed")+".pdf")
        def task():
            try:
                before=os.path.getsize(src)/1024
                self._log(f"Compressing {Path(src).name}  ({before:.1f} KB) …")
                compress_pdf(src,dst)
                after=os.path.getsize(dst)/1024
                self._log(f"✔ {Path(dst).name}  ({after:.1f} KB)  — saved {before-after:.1f} KB","ok")
                self.after(300,self._open_output)
            except Exception as ex: self._log(f"ERROR: {ex}","err")
        threading.Thread(target=task,daemon=True).start()

    def _run_metadata(self):
        src=self._require_pdf("Metadata"); out=self._require_out()
        if not src or not out: return
        dst=os.path.join(out,(self.meta_out.get().strip() or "updated")+".pdf")
        fields={k:v.get() for k,v in self.meta_fields.items()}
        def task():
            try:
                set_metadata(src,dst,fields)
                self._log(f"✔ {Path(dst).name}","ok")
                self.after(300,self._open_output)
            except Exception as ex: self._log(f"ERROR: {ex}","err")
        threading.Thread(target=task,daemon=True).start()

    # ─── File queue ───────────────────────────────────────────────────────────

    def _add_files(self):
        paths=filedialog.askopenfilenames(
            title="Select files",
            filetypes=[
                ("All supported",
                 "*.pdf *.docx *.xlsx *.xls *.xlsm *.pptx *.ppt *.csv *.tsv "
                 "*.txt *.html *.htm "
                 "*.png *.jpg *.jpeg *.webp *.bmp *.gif *.tiff *.ico "
                 "*.mp4 *.avi *.mov *.mkv *.webm *.mp3 *.wav *.ogg *.flac *.aac *.m4a"),
                ("PDF","*.pdf"),
                ("Excel","*.xlsx *.xls *.xlsm *.ods"),
                ("PowerPoint","*.pptx *.ppt *.odp"),
                ("CSV / TSV","*.csv *.tsv"),
                ("Word","*.docx"),
                ("Images","*.png *.jpg *.jpeg *.webp *.bmp *.gif *.tiff *.ico"),
                ("Audio","*.mp3 *.wav *.ogg *.flac *.aac *.m4a"),
                ("Video","*.mp4 *.avi *.mov *.mkv *.webm"),
                ("All files","*.*"),
            ])
        for p in paths:
            if p not in self.files:
                self.files.append(p); c=cat(p)
                self.queue_lb.insert("end",f" {cat_icon(c)}  {Path(p).name}  [{c}]")
        self._update_after_files()

    def _add_folder(self):
        folder=filedialog.askdirectory(title="Add all files from folder")
        if not folder: return
        exts=IMAGE_EXTS|AUDIO_EXTS|VIDEO_EXTS|EXCEL_EXTS|PPTX_EXTS|CSV_EXTS|{".pdf",".docx",".txt",".html",".htm"}
        added=0
        for path in sorted(Path(folder).iterdir()):
            if path.suffix.lower() in exts and str(path) not in self.files:
                self.files.append(str(path)); c=cat(str(path))
                self.queue_lb.insert("end",f" {cat_icon(c)}  {path.name}  [{c}]")
                added+=1
        self._log(f"Added {added} file(s) from folder.")
        self._update_after_files()

    def _update_after_files(self):
        n=len(self.files); self._status(f"{n} file(s) in queue")
        if not self.out_dir.get() and self.files:
            self.out_dir.set(str(Path(self.files[0]).parent/"workshop_output"))
        pdfs=[f for f in self.files if cat(f)=="pdf"]
        if pdfs and HAS_PYPDF:
            try:
                pc=pdf_page_count(pdfs[0]); self.page_count=pc; self.active_pdf=pdfs[0]
                name=Path(pdfs[0]).name
                self.pdf_badge.config(text=f"📄 {name}\n   {pc} pages")
                self.split_info_lbl.config(
                    text=f"Active PDF: {name}  ·  {pc} pages",fg=C["green"])
                self.org_pdf_lbl_var.set(f"Active PDF: {name}  ·  {pc} pages")
            except: pass
        pptxs=[f for f in self.files if cat(f)=="pptx"]
        if pptxs and HAS_PPTX:
            try:
                sc=pptx_slide_count(pptxs[0])
                self.split_info_lbl.config(
                    text=f"Active PPTX: {Path(pptxs[0]).name}  ·  {sc} slides",fg=C["green"])
            except: pass
        self._refresh_merge_preview()

    def _refresh_merge_preview(self):
        if not hasattr(self,"merge_preview"): return
        self.merge_preview.config(state="normal")
        self.merge_preview.delete("1.0","end")
        if not self.files:
            self.merge_preview.insert("end","  No files in queue.")
        else:
            for i,fp in enumerate(self.files,1):
                c=cat(fp)
                self.merge_preview.insert("end",f"  {i:2d}.  {cat_icon(c)}  {Path(fp).name}\n")
        self.merge_preview.config(state="disabled")

    def _move_up(self):
        sel=list(self.queue_lb.curselection())
        if not sel or sel[0]==0: return
        for i in sel:
            self.files[i-1],self.files[i]=self.files[i],self.files[i-1]
            a=self.queue_lb.get(i-1); b=self.queue_lb.get(i)
            self.queue_lb.delete(i-1,i)
            self.queue_lb.insert(i-1,b); self.queue_lb.insert(i,a)
        self.queue_lb.selection_clear(0,"end")
        for i in sel: self.queue_lb.selection_set(i-1)
        self._refresh_merge_preview()

    def _move_down(self):
        sel=list(self.queue_lb.curselection())
        if not sel or sel[-1]==len(self.files)-1: return
        for i in reversed(sel):
            self.files[i],self.files[i+1]=self.files[i+1],self.files[i]
            a=self.queue_lb.get(i); b=self.queue_lb.get(i+1)
            self.queue_lb.delete(i,i+1)
            self.queue_lb.insert(i,b); self.queue_lb.insert(i+1,a)
        self.queue_lb.selection_clear(0,"end")
        for i in sel: self.queue_lb.selection_set(i+1)
        self._refresh_merge_preview()

    def _remove_selected(self):
        sel=list(self.queue_lb.curselection())
        for i in reversed(sel): self.files.pop(i); self.queue_lb.delete(i)
        self._update_after_files()

    def _clear_queue(self):
        self.files.clear(); self.queue_lb.delete(0,"end")
        self.page_count=0; self.active_pdf=""
        self.pdf_badge.config(text=""); self._status("Queue cleared")
        self._refresh_merge_preview()

    def _browse_output(self):
        p=filedialog.askdirectory(title="Select output folder")
        if p: self.out_dir.set(p)

    def _browse_stamp_pdf(self):
        p=filedialog.askopenfilename(title="Select watermark PDF",filetypes=[("PDF","*.pdf")])
        if p: self.stamp_wm_path.set(p)

    def _open_output(self):
        folder=self.out_dir.get()
        if not folder: return
        os.makedirs(folder,exist_ok=True)
        try:
            if sys.platform=="win32": os.startfile(folder)
            elif sys.platform=="darwin": subprocess.run(["open",folder])
            else: subprocess.run(["xdg-open",folder])
        except Exception as e: self._log(f"Could not open folder: {e}","err")

    def _load_metadata(self):
        src=self._require_pdf("Metadata")
        if not src: return
        try:
            meta=get_metadata(src)
            self.meta_display.config(state="normal")
            self.meta_display.delete("1.0","end")
            for k,v in meta.items(): self.meta_display.insert("end",f"  {k:12s}: {v}\n")
            self.meta_display.config(state="disabled")
            mapping={"Title":"title","Author":"author","Subject":"subject","Creator":"creator"}
            for dk,fk in mapping.items():
                val=meta.get(dk,"")
                if val!="—" and fk in self.meta_fields: self.meta_fields[fk].set(val)
        except Exception as e: self._log(f"Could not load metadata: {e}","err")

    def _switch_tab(self,key:str):
        for k,btn in self.tab_btns.items():
            btn.config(bg=C["accent"] if k==key else C["sidebar"],
                       fg="#ffffff" if k==key else C["dim"])
        for k,page in self.pages.items(): page.pack_forget()
        self.pages[key].pack(fill="both",expand=True)
        if key=="merge": self._refresh_merge_preview()

    def _on_fmt_change(self,*_):
        fmt=self.out_fmt.get()
        warn={
            "docx":"" if HAS_DOCX else "⚠  pip install python-docx",
            "xlsx":"" if HAS_OPENPYXL else "⚠  pip install openpyxl",
            "csv": "",
            "pptx":"" if HAS_PPTX else "⚠  pip install python-pptx",
            "png": "" if HAS_PDF2IMAGE else "⚠  pip install pdf2image  +  sudo apt install poppler-utils",
            "jpg": "" if HAS_PDF2IMAGE else "⚠  pip install pdf2image  +  sudo apt install poppler-utils",
            "mp3": "" if HAS_FFMPEG else "⚠  ffmpeg not found — sudo apt install ffmpeg",
            "mp4": "" if HAS_FFMPEG else "⚠  ffmpeg not found",
            "gif": "" if HAS_FFMPEG else "⚠  ffmpeg not found",
        }
        lo_needed={"xlsx","pptx","csv"} if not HAS_LIBREOFFICE else set()
        note=warn.get(fmt,"")
        if fmt in ("xlsx","pptx") and not HAS_LIBREOFFICE:
            note+="  ⚠  LibreOffice needed for Office→PDF: sudo apt install libreoffice"
        self.conv_note.config(text=note)

    def _log(self,msg:str,kind:str=""):
        def _do():
            self.log_box.config(state="normal")
            tag={"ok":"ok","err":"err","info":"info","warn":"warn"}.get(kind,"")
            self.log_box.insert("end","  "+msg+"\n",tag)
            self.log_box.see("end"); self.log_box.config(state="disabled")
        self.after(0,_do)

    def _clear_log(self):
        self.log_box.config(state="normal"); self.log_box.delete("1.0","end")
        self.log_box.config(state="disabled")

    def _status(self,msg:str):
        self.after(0,lambda: self.status_var.set("  "+msg))

    def _refresh_deps(self):
        deps=[HAS_PYPDF,HAS_PDFPLUMBER,HAS_DOCX,HAS_PIL,HAS_PDF2IMAGE,
              HAS_FFMPEG,HAS_OPENPYXL,HAS_PPTX,HAS_LIBREOFFICE]
        n=sum(deps); total=len(deps)
        self.dep_lbl.config(text=f"deps {n}/{total} ✓",
                             fg=C["green"] if n==total else C["yellow"])
        Tooltip(self.dep_lbl,
            f"pypdf:        {'✔' if HAS_PYPDF else '✖  pip install pypdf'}\n"
            f"pdfplumber:   {'✔' if HAS_PDFPLUMBER else '✖  pip install pdfplumber'}\n"
            f"python-docx:  {'✔' if HAS_DOCX else '✖  pip install python-docx'}\n"
            f"Pillow:       {'✔' if HAS_PIL else '✖  pip install Pillow'}\n"
            f"pdf2image:    {'✔' if HAS_PDF2IMAGE else '✖  pip install pdf2image'}\n"
            f"ffmpeg:       {'✔' if HAS_FFMPEG else '✖  sudo apt install ffmpeg'}\n"
            f"openpyxl:     {'✔' if HAS_OPENPYXL else '✖  pip install openpyxl'}\n"
            f"python-pptx:  {'✔' if HAS_PPTX else '✖  pip install python-pptx'}\n"
            f"LibreOffice:  {'✔' if HAS_LIBREOFFICE else '✖  sudo apt install libreoffice'}")

    def _btn(self,parent,text,cmd,color=None):
        return tk.Button(parent,text=text,command=cmd,
                         bg=color or C["accent2"],fg="#ffffff",
                         activebackground=C["border"],relief="flat",
                         cursor="hand2",font=F["btn"],padx=10,pady=4,bd=0)

    def _bigbtn(self,parent,text,cmd,color=None):
        return tk.Button(parent,text=text,command=cmd,
                         bg=color or C["accent"],fg="#ffffff",
                         activebackground=C["border"],relief="flat",
                         cursor="hand2",font=F["btnbig"],padx=14,pady=8,bd=0)

    def _lbl(self,parent,text,head=False,dim=False,w=None,wrap=None):
        bg=C["panel"]
        try: bg=parent.cget("bg")
        except: pass
        kw=dict(bg=bg,fg=C["dim"] if dim else C["text"],
                font=F["head"] if head else F["label"],anchor="w")
        if w: kw["width"]=w
        if wrap: kw["wraplength"]=wrap; kw["justify"]="left"
        return tk.Label(parent,text=text,**kw)

    def _card(self,parent):
        f=tk.Frame(parent,bg=C["panel"],padx=16,pady=12)
        f.pack(fill="x",padx=14,pady=(0,6)); return f

    def _section_title(self,parent,title,sub=""):
        h=tk.Frame(parent,bg=C["bg"]); h.pack(fill="x",padx=14,pady=(12,4))
        tk.Label(h,text=title,font=F["head"],bg=C["bg"],fg=C["accent"]).pack(side="left")
        if sub: tk.Label(h,text=f"  —  {sub}",font=F["small"],bg=C["bg"],fg=C["dim"]).pack(side="left")


# ═══════════════════════════════════════════════════════════════════════════════
# HELP TEXT
# ═══════════════════════════════════════════════════════════════════════════════

HELP_TEXT = """FILE WORKSHOP v4 — COMPLETE REFERENCE
════════════════════════════════════════════════════════════════════

FULL CONVERSION TABLE
──────────────────────────────────────────────────────
  PDF      →  DOCX, TXT, HTML, PNG, JPG, XLSX, PPTX, CSV
  Excel    →  PDF*, CSV, TXT, HTML, DOCX
             (XLSX↔XLS↔ODS via LibreOffice)
  CSV/TSV  →  XLSX, PDF, HTML, TXT, DOCX
  PPTX     →  PDF*, TXT, HTML, DOCX, PNG, JPG
  DOCX     →  PDF*, TXT, HTML
  TXT      →  PDF, DOCX, HTML
  HTML     →  PDF**, TXT
  Images   →  PNG, JPG, WEBP, BMP, GIF, TIFF, ICO, PDF
  Video    →  MP4, AVI, MOV, MKV, WEBM, GIF, MP3, WAV  (ffmpeg)
  Audio    →  MP3, WAV, OGG, FLAC, AAC, M4A             (ffmpeg)

  *  PDF output from Office files needs LibreOffice
  ** HTML→PDF needs weasyprint

SPLIT
──────
  PDF:   Each page | Range | Custom groups
  PPTX:  Each slide → own .pptx file

MERGE
──────
  PDF:  Combine PDFs + images into one PDF
  PPTX: Combine multiple .pptx into one presentation

ORGANISE (PDF)
───────────────
  Resequence:  Enter new order e.g. 3,1,2
  Delete:      Remove pages e.g. 2,5,7-9
  Rotate:      90/180/270° on all or specific pages
  Reverse:     Flip entire page order

STAMP / WATERMARK (PDF)
─────────────────────────
  Text: diagonal text overlay (color, opacity, angle, font size)
  PDF overlay: merge a stamp PDF over target pages

PROTECT (PDF)
──────────────
  Encrypt: add user + optional owner password
  Decrypt: remove password (need current password)

COMPRESS (PDF)
───────────────
  Lossless stream compression — most effective on text-heavy PDFs

METADATA (PDF)
───────────────
  View Title/Author/Subject/Creator/Producer/Dates
  Edit and save back to a new PDF

PAGE/SLIDE SYNTAX
──────────────────
  1,3,5-8   →  pages 1, 3, 5, 6, 7, 8
  2-10      →  pages 2 through 10
  (blank)   →  all pages

CUSTOM GROUP SYNTAX (Split)
────────────────────────────
  1,10 | 3,5 | 2-4,7
  Creates 3 files: {1,10}  {3,5}  {2,3,4,7}

DEPENDENCIES
─────────────
  pip install pypdf pdfplumber python-docx reportlab Pillow pdf2image openpyxl python-pptx weasyprint

  sudo apt install libreoffice      # Office→PDF, XLSX↔format
  sudo apt install ffmpeg           # audio/video
  sudo apt install poppler-utils    # PDF→image, watermark
"""

# ═══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    if not HAS_PYPDF:
        print("pypdf is required:  pip install pypdf"); sys.exit(1)
    App().mainloop()
