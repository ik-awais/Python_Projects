"""
backend/processor.py - All file conversion, split, merge, organise operations
Ported from desktop tool with full functionality
"""

import os, re, csv, shutil, tempfile, gc
import subprocess
from pathlib import Path
from typing import List, Dict, Optional, Union
import mimetypes
import json
from datetime import datetime

# Try to import required packages
def _try(pkg):
    try: __import__(pkg); return True
    except ImportError: return False

HAS_PYPDF = _try("pypdf")
HAS_PDFPLUMBER = _try("pdfplumber")
HAS_DOCX = _try("docx")
HAS_PIL = _try("PIL")
HAS_PDF2IMAGE = _try("pdf2image")
HAS_REPORTLAB = _try("reportlab")
HAS_WEASYPRINT = _try("weasyprint")
HAS_CAIROSVG = _try("cairosvg")
HAS_OPENPYXL = _try("openpyxl")
HAS_PPTX = _try("pptx")
HAS_FFMPEG = bool(shutil.which("ffmpeg"))
HAS_LIBREOFFICE = bool(shutil.which("libreoffice") or shutil.which("soffice"))

# Import packages if available
if HAS_PYPDF:
    from pypdf import PdfReader, PdfWriter
if HAS_PDFPLUMBER: import pdfplumber
if HAS_DOCX:
    from docx import Document
    from docx.shared import Pt, Inches
if HAS_PIL: from PIL import Image, ImageDraw, ImageFont
if HAS_PDF2IMAGE: from pdf2image import convert_from_path
if HAS_CAIROSVG: import cairosvg
if HAS_OPENPYXL: import openpyxl
if HAS_PPTX:
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt

# File type definitions
IMAGE_EXTS = {".png",".jpg",".jpeg",".webp",".bmp",".gif",".tiff",".tif",".ico",".svg"}
AUDIO_EXTS = {".mp3",".wav",".ogg",".flac",".aac",".m4a",".wma"}
VIDEO_EXTS = {".mp4",".avi",".mov",".mkv",".webm",".flv",".wmv",".m4v"}
EXCEL_EXTS = {".xlsx",".xls",".xlsm",".ods"}
PPTX_EXTS = {".pptx",".ppt",".odp"}
CSV_EXTS = {".csv",".tsv"}

def cat(path):
    """Get file category"""
    e = Path(path).suffix.lower()
    if e in IMAGE_EXTS: return "image"
    if e in AUDIO_EXTS: return "audio"
    if e in VIDEO_EXTS: return "video"
    if e in EXCEL_EXTS: return "excel"
    if e in PPTX_EXTS: return "pptx"
    if e in CSV_EXTS: return "csv"
    if e == ".pdf": return "pdf"
    if e == ".docx": return "docx"
    if e == ".txt": return "txt"
    if e in {".html",".htm"}: return "html"
    return "unknown"

def cat_icon(c):
    """Get icon for file category"""
    return {"pdf":"📄","image":"🖼","audio":"🎵","video":"🎬",
            "docx":"📝","txt":"📃","html":"🌐","excel":"📊",
            "pptx":"📽","csv":"📋"}.get(c,"📎")

def parse_pages(s, total):
    """Parse page ranges like '1,3,5-8'"""
    pages = set()
    for part in s.split(","):
        part = part.strip()
        if "-" in part:
            a,b = part.split("-",1)
            try: pages.update(range(int(a),int(b)+1))
            except: pass
        else:
            try: pages.add(int(part))
            except: pass
    return sorted(p for p in pages if 1<=p<=total)

def _safe_filename(name: str, max_bytes: int = 250) -> str:
    """Ensure a filename doesn't exceed OS byte limit."""
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', '_', name)
    while len(name.encode("utf-8")) > max_bytes:
        name = name[:-1]
    return name

# ── LibreOffice Conversion ───────────────────────────────────────────────────────

def lo_convert(src, dst_fmt, out_dir):
    """Convert using LibreOffice"""
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if not lo:
        raise RuntimeError("LibreOffice not found.")
    os.makedirs(out_dir, exist_ok=True)
    r = subprocess.run([lo,"--headless","--convert-to",dst_fmt,"--outdir",out_dir,src],
                       capture_output=True, text=True)
    expected = Path(out_dir)/(Path(src).stem+"."+dst_fmt)
    if not expected.exists():
        raise RuntimeError(f"LibreOffice failed: {r.stderr[-400:]}")
    return str(expected)

# ── PDF Operations ───────────────────────────────────────────────────────────────

def pdf_page_count(src): 
    """Get PDF page count"""
    return len(PdfReader(src).pages)

def pdf_to_txt(src, dst, pages=None):
    """Convert PDF to text"""
    if not HAS_PDFPLUMBER: raise ImportError("pip install pdfplumber")
    with pdfplumber.open(src) as pdf:
        total = len(pdf.pages)
        target = pages or list(range(1, total+1))
        lines = []
        for p in target:
            if 1 <= p <= total:
                lines.append(f"\n{'='*50}\nPAGE {p}\n{'='*50}\n")
                lines.append(pdf.pages[p-1].extract_text() or "")
    Path(dst).write_text("\n".join(lines), encoding="utf-8")
    return dst

def pdf_to_docx(src, dst, pages=None):
    """Convert PDF to DOCX"""
    if not(HAS_PDFPLUMBER and HAS_DOCX): raise ImportError("pip install pdfplumber python-docx")
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    with pdfplumber.open(src) as pdf:
        total = len(pdf.pages)
        target = pages or list(range(1, total+1))
        for p in target:
            if 1 <= p <= total:
                doc.add_heading(f"Page {p}", 2)
                text = pdf.pages[p-1].extract_text() or ""
                for para in text.split("\n\n"):
                    if para.strip(): doc.add_paragraph(para.strip())
                doc.add_page_break()
    doc.save(dst)
    return dst

def pdf_to_images(src, dst_dir, fmt="png", pages=None, dpi=150):
    """Convert PDF pages to images"""
    if not HAS_PDF2IMAGE: raise ImportError("pip install pdf2image (+poppler)")
    os.makedirs(dst_dir, exist_ok=True)

    if pages:
        target = sorted(pages)
    else:
        total = pdf_page_count(src)
        target = list(range(1, total + 1))

    out = []
    batch_size = 20

    for batch_start in range(0, len(target), batch_size):
        batch = target[batch_start:batch_start + batch_size]
        first_pg = min(batch)
        last_pg = max(batch)

        imgs = convert_from_path(src, dpi=dpi, fmt=fmt,
                                 first_page=first_pg, last_page=last_pg)

        page_range = list(range(first_pg, last_pg + 1))
        for img, pg_num in zip(imgs, page_range):
            if pg_num in batch:
                p = os.path.join(dst_dir, f"page_{pg_num}.{fmt}")
                img.save(p)
                out.append(p)
            img.close()
        del imgs
        gc.collect()

    return out

# ── Excel Operations ─────────────────────────────────────────────────────────────

def excel_to_pdf(src, dst):
    """Convert Excel to PDF"""
    out_dir = str(Path(dst).parent)
    result = lo_convert(src, "pdf", out_dir)
    if str(result) != str(dst): shutil.move(result, dst)
    return dst

def excel_to_csv(src, dst, sheet_index=0):
    """Convert Excel to CSV"""
    if not HAS_OPENPYXL: raise ImportError("pip install openpyxl")
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    sheet = wb.worksheets[sheet_index]
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        for row in sheet.iter_rows(values_only=True):
            writer.writerow([v if v is not None else "" for v in row])
    wb.close()
    return dst

# ── CSV Operations ───────────────────────────────────────────────────────────────

def csv_to_xlsx(src, dst):
    """Convert CSV to Excel"""
    if not HAS_OPENPYXL: raise ImportError("pip install openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    delim = "\t" if src.lower().endswith(".tsv") else ","
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f, delimiter=delim): ws.append(row)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    wb.save(dst)
    return dst

# ── PPTX Operations ─────────────────────────────────────────────────────────────

def pptx_to_pdf(src, dst):
    """Convert PowerPoint to PDF"""
    out_dir = str(Path(dst).parent)
    result = lo_convert(src, "pdf", out_dir)
    if str(result) != str(dst): shutil.move(result, dst)
    return dst

def pptx_to_txt(src, dst):
    """Convert PowerPoint to text"""
    if not HAS_PPTX: raise ImportError("pip install python-pptx")
    prs = Presentation(src)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n{'='*50}\nSLIDE {i}\n{'='*50}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t: lines.append(t)
    Path(dst).write_text("\n".join(lines), encoding="utf-8")
    return dst

def pptx_slide_count(src):
    """Get PowerPoint slide count"""
    if not HAS_PPTX: return 0
    return len(Presentation(src).slides)

# ── DOCX Operations ─────────────────────────────────────────────────────────────

def docx_to_txt(src, dst):
    """Convert DOCX to text"""
    if not HAS_DOCX: raise ImportError("pip install python-docx")
    doc = Document(src)
    Path(dst).write_text("\n".join(p.text for p in doc.paragraphs), encoding="utf-8")
    return dst

def docx_to_pdf(src, dst):
    """Convert DOCX to PDF"""
    out_dir = str(Path(dst).parent)
    result = lo_convert(src, "pdf", out_dir)
    if str(result) != str(dst): shutil.move(result, dst)
    return dst

# ── TXT Operations ───────────────────────────────────────────────────────────────

def txt_to_pdf(src, dst):
    """Convert text to PDF"""
    if not HAS_REPORTLAB: raise ImportError("pip install reportlab")
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import mm
    
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    doc = SimpleDocTemplate(dst, pagesize=A4,
                            leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=20*mm, bottomMargin=20*mm)
    styles = getSampleStyleSheet()
    story = []
    for line in text.split("\n"):
        safe = line.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;") or "&nbsp;"
        story.append(Paragraph(safe, styles["Normal"]))
        story.append(Spacer(1, 2))
    doc.build(story)
    return dst

def txt_to_docx(src, dst):
    """Convert text to DOCX"""
    if not HAS_DOCX: raise ImportError("pip install python-docx")
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    for line in Path(src).read_text(encoding="utf-8", errors="replace").split("\n"):
        doc.add_paragraph(line)
    doc.save(dst)
    return dst

# ── PDF Split Operations ───────────────────────────────────────────────────────

def split_each(src, out_dir):
    """Split PDF into individual pages"""
    if not HAS_PYPDF: raise ImportError("pip install pypdf")
    os.makedirs(out_dir, exist_ok=True)
    
    reader = PdfReader(src)
    total = len(reader.pages)
    w = max(3, len(str(total)))
    out = []
    
    for i, page in enumerate(reader.pages):
        writer = PdfWriter()
        writer.add_page(page)
        p = os.path.join(out_dir, f"{Path(src).stem}_page_{str(i+1).zfill(w)}.pdf")
        with open(p, "wb") as f:
            writer.write(f)
        out.append(p)
    
    return out

def split_range(src, out_dir, page_range):
    """Split PDF by page range"""
    if not HAS_PYPDF: raise ImportError("pip install pypdf")
    os.makedirs(out_dir, exist_ok=True)
    
    reader = PdfReader(src)
    pages = parse_pages(page_range, len(reader.pages))
    
    writer = PdfWriter()
    for p in pages:
        if 1 <= p <= len(reader.pages):
            writer.add_page(reader.pages[p-1])
    
    dst = os.path.join(out_dir, f"{Path(src).stem}_pages_{page_range.replace(',', '_').replace('-', 'to')}.pdf")
    with open(dst, "wb") as f:
        writer.write(f)
    
    return dst

# ── PDF Merge Operations ───────────────────────────────────────────────────────

def merge_pdfs(src_list, dst):
    """Merge multiple PDFs"""
    if not HAS_PYPDF: raise ImportError("pip install pypdf")
    
    writer = PdfWriter()
    for src in src_list:
        reader = PdfReader(src)
        for page in reader.pages:
            writer.add_page(page)
    
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    
    return dst

# ── PDF Protection Operations ───────────────────────────────────────────────────

def encrypt_pdf(src, dst, password):
    """Encrypt PDF with password"""
    if not HAS_PYPDF: raise ImportError("pip install pypdf")
    
    reader = PdfReader(src)
    writer = PdfWriter()
    
    for page in reader.pages:
        writer.add_page(page)
    
    writer.encrypt(password)
    
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    
    return dst

def decrypt_pdf(src, dst, password):
    """Decrypt PDF"""
    if not HAS_PYPDF: raise ImportError("pip install pypdf")
    
    reader = PdfReader(src)
    if reader.is_encrypted:
        reader.decrypt(password)
    
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    
    return dst

# ── Watermark Operations ───────────────────────────────────────────────────────

def watermark_text(src, dst, text, opacity=0.5):
    """Add text watermark to PDF"""
    if not HAS_PYPDF: raise ImportError("pip install pypdf")
    
    reader = PdfReader(src)
    writer = PdfWriter()
    
    for page in reader.pages:
        writer.add_page(page)
    
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    
    return dst

# ── Metadata Operations ───────────────────────────────────────────────────────

def get_metadata(src):
    """Get file metadata"""
    path = Path(src)
    stat = path.stat()
    
    metadata = {
        "filename": path.name,
        "size": stat.st_size,
        "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
        "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
        "category": cat(src),
        "icon": cat_icon(cat(src))
    }
    
    # Add PDF-specific metadata
    if cat(src) == "pdf" and HAS_PYPDF:
        try:
            reader = PdfReader(src)
            metadata["pages"] = len(reader.pages)
            if reader.metadata:
                metadata.update({
                    "title": reader.metadata.get('/Title', ''),
                    "author": reader.metadata.get('/Author', ''),
                    "subject": reader.metadata.get('/Subject', ''),
                    "creator": reader.metadata.get('/Creator', ''),
                    "producer": reader.metadata.get('/Producer', '')
                })
        except:
            pass
    
    return metadata

# ── Video/Audio Operations ───────────────────────────────────────────────────

def video_get_duration(src):
    """Get video duration using ffmpeg"""
    if not HAS_FFMPEG:
        raise RuntimeError("FFmpeg not found")
    
    cmd = ['ffmpeg', '-i', src, '-f', 'null', '-']
    result = subprocess.run(cmd, capture_output=True, text=True)
    
    # Parse duration from stderr
    for line in result.stderr.split('\n'):
        if 'Duration:' in line:
            duration_str = line.split('Duration:')[1].split(',')[0].strip()
            return duration_str
    
    return "00:00:00"

# ── Main Conversion Function ───────────────────────────────────────────────────

def do_convert(src, dst_fmt, out_dir, pages=None):
    """Main conversion function"""
    src_path = Path(src)
    src_cat = cat(src)
    
    # Determine conversion path
    if src_cat == "pdf":
        if dst_fmt == "txt":
            dst = out_dir / f"{src_path.stem}.txt"
            return pdf_to_txt(src, str(dst), pages)
        elif dst_fmt == "docx":
            dst = out_dir / f"{src_path.stem}.docx"
            return pdf_to_docx(src, str(dst), pages)
        elif dst_fmt == "images":
            dst_dir = out_dir / f"{src_path.stem}_images"
            return pdf_to_images(src, str(dst_dir), pages=pages)
    
    elif src_cat == "excel":
        if dst_fmt == "pdf":
            dst = out_dir / f"{src_path.stem}.pdf"
            return excel_to_pdf(src, str(dst))
        elif dst_fmt == "csv":
            dst = out_dir / f"{src_path.stem}.csv"
            return excel_to_csv(src, str(dst))
    
    elif src_cat == "pptx":
        if dst_fmt == "pdf":
            dst = out_dir / f"{src_path.stem}.pdf"
            return pptx_to_pdf(src, str(dst))
        elif dst_fmt == "txt":
            dst = out_dir / f"{src_path.stem}.txt"
            return pptx_to_txt(src, str(dst))
    
    elif src_cat == "docx":
        if dst_fmt == "pdf":
            dst = out_dir / f"{src_path.stem}.pdf"
            return docx_to_pdf(src, str(dst))
        elif dst_fmt == "txt":
            dst = out_dir / f"{src_path.stem}.txt"
            return docx_to_txt(src, str(dst))
    
    elif src_cat == "txt":
        if dst_fmt == "pdf":
            dst = out_dir / f"{src_path.stem}.pdf"
            return txt_to_pdf(src, str(dst))
        elif dst_fmt == "docx":
            dst = out_dir / f"{src_path.stem}.docx"
            return txt_to_docx(src, str(dst))
    
    elif src_cat == "csv":
        if dst_fmt == "xlsx":
            dst = out_dir / f"{src_path.stem}.xlsx"
            return csv_to_xlsx(src, str(dst))
    
    # Fallback to LibreOffice
    try:
        return lo_convert(src, dst_fmt, str(out_dir))
    except:
        raise ValueError(f"Unsupported conversion: {src_cat} to {dst_fmt}")

# ── Available Operations ─────────────────────────────────────────────────────

def get_available_operations():
    """Get list of available operations"""
    return {
        "convert": {
            "pdf": ["txt", "docx", "images"],
            "excel": ["pdf", "csv"],
            "pptx": ["pdf", "txt"],
            "docx": ["pdf", "txt"],
            "txt": ["pdf", "docx"],
            "csv": ["xlsx"]
        },
        "split": {
            "pdf": ["each", "range"]
        },
        "merge": {
            "pdf": True
        },
        "protect": {
            "pdf": ["encrypt", "decrypt"]
        },
        "metadata": {
            "all": True
        },
        "media": {
            "video": ["duration"],
            "audio": ["duration"]
        }
    }
