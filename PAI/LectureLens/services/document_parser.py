"""Unified document parser for PDF, DOCX, PPTX."""
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any

logger = logging.getLogger(__name__)

# PDF
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# DOCX
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

# PPTX
try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def parse_document(file_path: Path) -> Dict[str, Any]:
    """
    Extract text and metadata from a document.

    Returns:
        {
            "metadata": {"filename": str, "page_count": int, "file_size": int},
            "pages": [{"page_num": int, "text": str}, ...]
        }
    Raises:
        ValueError: Unsupported or corrupted/encrypted file.
    """
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = file_path.suffix.lower()
    file_size = file_path.stat().st_size

    if ext == '.pdf':
        return _parse_pdf(file_path, file_size)
    elif ext == '.docx':
        return _parse_docx(file_path, file_size)
    elif ext == '.pptx':
        return _parse_pptx(file_path, file_size)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _parse_pdf(file_path: Path, file_size: int) -> Dict[str, Any]:
    if fitz is None:
        raise ImportError("PyMuPDF (fitz) not installed")

    try:
        doc = fitz.open(file_path)
    except Exception as e:
        raise ValueError(f"Failed to open PDF (possibly encrypted or corrupted): {e}")

    pages = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        pages.append({"page_num": page_num + 1, "text": text.strip() or ""})

    metadata = {
        "filename": file_path.name,
        "page_count": len(pages),
        "file_size": file_size,
    }
    doc.close()
    return {"metadata": metadata, "pages": pages}


def _parse_docx(file_path: Path, file_size: int) -> Dict[str, Any]:
    if DocxDocument is None:
        raise ImportError("python-docx not installed")

    try:
        doc = DocxDocument(file_path)
    except Exception as e:
        raise ValueError(f"Failed to open DOCX: {e}")

    # Treat each paragraph as a page? DOCX has no natural pages.
    # We'll combine all paragraphs into one "page" (page_num=1).
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)
    text = "\n".join(full_text).strip()

    pages = [{"page_num": 1, "text": text}] if text else []

    metadata = {
        "filename": file_path.name,
        "page_count": len(pages),
        "file_size": file_size,
    }
    return {"metadata": metadata, "pages": pages}


def _parse_pptx(file_path: Path, file_size: int) -> Dict[str, Any]:
    if Presentation is None:
        raise ImportError("python-pptx not installed")

    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f"Failed to open PPTX: {e}")

    pages = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        text = "\n".join(slide_text).strip()
        if text:  # skip empty slides
            pages.append({"page_num": slide_num, "text": text})

    metadata = {
        "filename": file_path.name,
        "page_count": len(pages),
        "file_size": file_size,
    }
    return {"metadata": metadata, "pages": pages}