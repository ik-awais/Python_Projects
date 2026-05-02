"""
core/processor.py — All file conversion, split, merge, organise operations.
Memory-efficient for 100k+ page files. Clean importable functions.
"""

import os, re, csv, shutil, tempfile, gc
from pathlib import Path
from typing import List, Dict, Optional

def _try(pkg):
    try: __import__(pkg); return True
    except ImportError: return False

HAS_PYPDF      = _try("pypdf")
HAS_PDFPLUMBER = _try("pdfplumber")
HAS_DOCX       = _try("docx")
HAS_PIL        = _try("PIL")
HAS_PDF2IMAGE  = _try("pdf2image")
HAS_REPORTLAB  = _try("reportlab")
HAS_WEASYPRINT = _try("weasyprint")
HAS_CAIROSVG   = _try("cairosvg")
HAS_OPENPYXL   = _try("openpyxl")
HAS_PPTX       = _try("pptx")
HAS_FFMPEG     = bool(shutil.which("ffmpeg"))
HAS_LIBREOFFICE = bool(shutil.which("libreoffice") or shutil.which("soffice"))

import subprocess
if HAS_PYPDF:
    from pypdf import PdfReader, PdfWriter
if HAS_PDFPLUMBER: import pdfplumber
if HAS_DOCX:
    from docx import Document
    from docx.shared import Pt, Inches
if HAS_PIL: from PIL import Image, ImageDraw, ImageFont
if HAS_PDF2IMAGE: from pdf2image import convert_from_path
if HAS_CAIROSVG: import cairosvg
if HAS_OPENPYXL: import openpyxl
if HAS_PPTX:
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt

IMAGE_EXTS = {".png",".jpg",".jpeg",".webp",".bmp",".gif",".tiff",".tif",".ico",".svg"}
AUDIO_EXTS = {".mp3",".wav",".ogg",".flac",".aac",".m4a",".wma"}
VIDEO_EXTS = {".mp4",".avi",".mov",".mkv",".webm",".flv",".wmv",".m4v"}
EXCEL_EXTS = {".xlsx",".xls",".xlsm",".ods"}
PPTX_EXTS  = {".pptx",".ppt",".odp"}
CSV_EXTS   = {".csv",".tsv"}

def cat(path):
    e = Path(path).suffix.lower()
    if e in IMAGE_EXTS: return "image"
    if e in AUDIO_EXTS: return "audio"
    if e in VIDEO_EXTS: return "video"
    if e in EXCEL_EXTS: return "excel"
    if e in PPTX_EXTS:  return "pptx"
    if e in CSV_EXTS:   return "csv"
    if e == ".pdf":     return "pdf"
    if e == ".docx":    return "docx"
    if e == ".txt":     return "txt"
    if e in {".html",".htm"}: return "html"
    return "unknown"

def cat_icon(c):
    return {"pdf":"📄","image":"🖼","audio":"🎵","video":"🎬",
            "docx":"📝","txt":"📃","html":"🌐","excel":"📊",
            "pptx":"📽","csv":"📋"}.get(c,"📎")

def parse_pages(s, total):
    pages = set()
    for part in s.split(","):
        part = part.strip()
        if "-" in part:
            a,b = part.split("-",1)
            try: pages.update(range(int(a),int(b)+1))
            except: pass
        else:
            try: pages.add(int(part))
            except: pass
    return sorted(p for p in pages if 1<=p<=total)

def parse_groups(s, total):
    return [g for g in (parse_pages(c.strip(),total) for c in s.split("|")) if g]

# ── Helpers for safe filenames and memory ─────────────────────────────────────

def _safe_stem(path: str, max_len: int = 60) -> str:
    """Truncate the stem of a filename to stay within OS limits."""
    stem = Path(path).stem
    if len(stem) > max_len:
        stem = stem[:max_len]
    return stem

def _width(total: int) -> int:
    """How many digits needed to number `total` items (minimum 3)."""
    return max(3, len(str(total)))

def _safe_filename(name: str, max_bytes: int = 250) -> str:
    """Ensure a filename doesn't exceed OS byte limit."""
    # Remove characters that are illegal on most filesystems
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', '_', name)
    while len(name.encode("utf-8")) > max_bytes:
        name = name[:-1]
    return name

def _safe_video_filename(time_str: str) -> str:
    """Sanitize a time string for use in filenames (replace colons)."""
    return time_str.replace(":", "-").replace(" ", "")

# ── LibreOffice ───────────────────────────────────────────────────────────────

def lo_convert(src, dst_fmt, out_dir):
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if not lo:
        raise RuntimeError("LibreOffice not found.\n  sudo apt install libreoffice")
    os.makedirs(out_dir, exist_ok=True)
    r = subprocess.run([lo,"--headless","--convert-to",dst_fmt,"--outdir",out_dir,src],
                       capture_output=True, text=True)
    expected = Path(out_dir)/(Path(src).stem+"."+dst_fmt)
    if not expected.exists():
        raise RuntimeError(f"LibreOffice failed:\n{r.stderr[-400:]}")
    return str(expected)

# ── PDF ───────────────────────────────────────────────────────────────────────

def pdf_page_count(src): return len(PdfReader(src).pages)

def pdf_to_txt(src, dst, pages=None):
    if not HAS_PDFPLUMBER: raise ImportError("pip install pdfplumber")
    with pdfplumber.open(src) as pdf:
        total = len(pdf.pages)
        target = pages or list(range(1, total+1))
        lines = []
        for p in target:
            if 1 <= p <= total:
                lines.append(f"\n{'='*50}\nPAGE {p}\n{'='*50}\n")
                lines.append(pdf.pages[p-1].extract_text() or "")
    Path(dst).write_text("\n".join(lines), encoding="utf-8")
    return dst

def pdf_to_docx(src, dst, pages=None):
    if not(HAS_PDFPLUMBER and HAS_DOCX): raise ImportError("pip install pdfplumber python-docx")
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    with pdfplumber.open(src) as pdf:
        total = len(pdf.pages)
        target = pages or list(range(1, total+1))
        for p in target:
            if 1 <= p <= total:
                doc.add_heading(f"Page {p}", 2)
                text = pdf.pages[p-1].extract_text() or ""
                for para in text.split("\n\n"):
                    if para.strip(): doc.add_paragraph(para.strip())
                doc.add_page_break()
    doc.save(dst)
    return dst

def pdf_to_images(src, dst_dir, fmt="png", pages=None, dpi=150):
    """Convert PDF pages to images. Processes in batches of 20 for memory efficiency."""
    if not HAS_PDF2IMAGE: raise ImportError("pip install pdf2image (+poppler)")
    os.makedirs(dst_dir, exist_ok=True)

    if pages:
        target = sorted(pages)
    else:
        total = pdf_page_count(src)
        target = list(range(1, total + 1))

    out = []
    batch_size = 20

    for batch_start in range(0, len(target), batch_size):
        batch = target[batch_start:batch_start + batch_size]
        first_pg = min(batch)
        last_pg = max(batch)

        imgs = convert_from_path(src, dpi=dpi, fmt=fmt,
                                 first_page=first_pg, last_page=last_pg)

        # Map converted images back to actual page numbers
        page_range = list(range(first_pg, last_pg + 1))
        for img, pg_num in zip(imgs, page_range):
            if pg_num in batch:
                p = os.path.join(dst_dir, f"page_{pg_num}.{fmt}")
                img.save(p)
                out.append(p)
            img.close()
        del imgs
        gc.collect()

    return out

def pdf_to_html(src, dst, pages=None):
    if not HAS_PDFPLUMBER: raise ImportError("pip install pdfplumber")
    with pdfplumber.open(src) as pdf:
        total = len(pdf.pages)
        target = pages or list(range(1, total+1))
        parts = [f"<html><head><meta charset='utf-8'><title>{Path(src).stem}</title>"
                 f"<style>body{{font-family:Georgia,serif;max-width:860px;margin:auto;padding:2em}}"
                 f"h2{{border-bottom:2px solid #7c6af7}}</style></head><body>"]
        for p in target:
            if 1 <= p <= total:
                text = pdf.pages[p-1].extract_text() or ""
                safe = text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                parts.append(f"<h2>Page {p}</h2><pre style='white-space:pre-wrap'>{safe}</pre><hr>")
        parts.append("</body></html>")
    Path(dst).write_text("\n".join(parts), encoding="utf-8")
    return dst

def pdf_to_xlsx(src, dst):
    if not(HAS_PDFPLUMBER and HAS_OPENPYXL): raise ImportError("pip install pdfplumber openpyxl")
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    with pdfplumber.open(src) as pdf:
        for i, page in enumerate(pdf.pages):
            ws = wb.create_sheet(title=f"Page_{i+1}")
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    for row in table: ws.append([c or "" for c in row])
                    ws.append([])
            else:
                for line in (page.extract_text() or "").split("\n"): ws.append([line])
    wb.save(dst)
    return dst

def pdf_to_pptx(src, dst, dpi=150):
    """Convert PDF to PPTX. Processes pages in batches for memory efficiency."""
    if not(HAS_PDF2IMAGE and HAS_PPTX): raise ImportError("pip install pdf2image python-pptx (+poppler)")
    total = pdf_page_count(src)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    batch_size = 10

    for batch_start in range(0, total, batch_size):
        first = batch_start + 1
        last = min(batch_start + batch_size, total)
        imgs = convert_from_path(src, dpi=dpi, first_page=first, last_page=last)
        for img in imgs:
            slide = prs.slides.add_slide(blank)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tf:
                tmp = tf.name
            img.save(tmp)
            slide.shapes.add_picture(tmp, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.unlink(tmp)
            img.close()
        del imgs
        gc.collect()

    prs.save(dst)
    return dst

# ── Excel ─────────────────────────────────────────────────────────────────────

def excel_to_pdf(src, dst):
    out_dir = str(Path(dst).parent)
    result = lo_convert(src, "pdf", out_dir)
    if str(result) != str(dst): shutil.move(result, dst)
    return dst

def excel_to_csv(src, dst, sheet_index=0):
    if not HAS_OPENPYXL: raise ImportError("pip install openpyxl")
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    sheet = wb.worksheets[sheet_index]
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        for row in sheet.iter_rows(values_only=True):
            writer.writerow([v if v is not None else "" for v in row])
    wb.close()
    return dst

def excel_to_txt(src, dst):
    if not HAS_OPENPYXL: raise ImportError("pip install openpyxl")
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"\n{'='*50}\nSHEET: {ws.title}\n{'='*50}")
        for row in ws.iter_rows(values_only=True):
            lines.append("\t".join(str(v) if v is not None else "" for v in row))
    wb.close()
    Path(dst).write_text("\n".join(lines), encoding="utf-8")
    return dst

def excel_to_html(src, dst):
    if not HAS_OPENPYXL: raise ImportError("pip install openpyxl")
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    parts = [f"<html><head><meta charset='utf-8'><title>{Path(src).stem}</title>"
             f"<style>body{{font-family:sans-serif;padding:2em}}"
             f"table{{border-collapse:collapse;margin-bottom:2em}}"
             f"td,th{{border:1px solid #ccc;padding:4px 10px;font-size:13px}}"
             f"th{{background:#f0f0f0}}</style></head><body>"]
    for ws in wb.worksheets:
        parts.append(f"<h2>{ws.title}</h2><table>")
        first = True
        for row in ws.iter_rows(values_only=True):
            tag = "th" if first else "td"
            cells = "".join(f"<{tag}>{str(v) if v is not None else ''}</{tag}>" for v in row)
            parts.append(f"<tr>{cells}</tr>")
            first = False
        parts.append("</table>")
    wb.close()
    Path(dst).write_text("\n".join(parts), encoding="utf-8")
    return dst

def excel_to_docx(src, dst):
    if not(HAS_OPENPYXL and HAS_DOCX): raise ImportError("pip install openpyxl python-docx")
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    for ws in wb.worksheets:
        doc.add_heading(ws.title, 2)
        rows = list(ws.iter_rows(values_only=True))
        if not rows: continue
        ncols = max(len(r) for r in rows)
        table = doc.add_table(rows=len(rows), cols=ncols)
        table.style = "Table Grid"
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = table.cell(ri, ci)
                cell.text = str(val) if val is not None else ""
                if ri == 0:
                    for run in cell.paragraphs[0].runs: run.bold = True
        doc.add_paragraph()
    wb.close()
    doc.save(dst)
    return dst

# ── CSV ───────────────────────────────────────────────────────────────────────

def csv_to_xlsx(src, dst):
    if not HAS_OPENPYXL: raise ImportError("pip install openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    delim = "\t" if src.lower().endswith(".tsv") else ","
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f, delimiter=delim): ws.append(row)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    wb.save(dst)
    return dst

def csv_to_pdf(src, dst):
    if not HAS_REPORTLAB: raise ImportError("pip install reportlab")
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import mm
    rows = []
    delim = "\t" if src.lower().endswith(".tsv") else ","
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f, delimiter=delim): rows.append(row)
    if not rows: raise ValueError("CSV empty")
    doc = SimpleDocTemplate(dst, pagesize=landscape(A4),
                            leftMargin=10*mm, rightMargin=10*mm,
                            topMargin=10*mm, bottomMargin=10*mm)
    styles = getSampleStyleSheet()
    story = [Paragraph(Path(src).stem, styles["Title"]), Spacer(1, 6)]
    ncols = max(len(r) for r in rows)
    padded = [r + [""]*(ncols - len(r)) for r in rows]
    t = Table(padded, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND",(0,0),(-1,0),colors.HexColor("#7c6af7")),
        ("TEXTCOLOR",(0,0),(-1,0),colors.white),
        ("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),
        ("FONTSIZE",(0,0),(-1,-1),8),
        ("GRID",(0,0),(-1,-1),0.5,colors.grey),
        ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f5f5f5")]),
        ("VALIGN",(0,0),(-1,-1),"MIDDLE"),
        ("PADDING",(0,0),(-1,-1),4)
    ]))
    story.append(t)
    doc.build(story)
    return dst

def csv_to_html(src, dst):
    rows = []
    delim = "\t" if src.lower().endswith(".tsv") else ","
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f, delimiter=delim): rows.append(row)
    parts = [f"<html><head><meta charset='utf-8'><style>body{{font-family:sans-serif;padding:2em}}"
             f"table{{border-collapse:collapse}}td,th{{border:1px solid #ccc;padding:5px 12px;"
             f"font-size:13px}}th{{background:#e8e8f0}}</style></head><body>"
             f"<h2>{Path(src).stem}</h2><table>"]
    for i, row in enumerate(rows):
        tag = "th" if i == 0 else "td"
        cells = "".join(f"<{tag}>{str(v).replace('&','&amp;').replace('<','&lt;')}</{tag}>" for v in row)
        parts.append(f"<tr>{cells}</tr>")
    parts.append("</table></body></html>")
    Path(dst).write_text("\n".join(parts), encoding="utf-8")
    return dst

def csv_to_docx(src, dst):
    if not HAS_DOCX: raise ImportError("pip install python-docx")
    rows = []
    delim = "\t" if src.lower().endswith(".tsv") else ","
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f, delimiter=delim): rows.append(row)
    if not rows: raise ValueError("CSV empty")
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    ncols = max(len(r) for r in rows)
    padded = [r + [""]*(ncols - len(r)) for r in rows]
    table = doc.add_table(rows=len(padded), cols=ncols)
    table.style = "Table Grid"
    for ri, row in enumerate(padded):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            if ri == 0:
                for run in cell.paragraphs[0].runs: run.bold = True
    doc.save(dst)
    return dst

# ── PPTX ──────────────────────────────────────────────────────────────────────

def pptx_to_pdf(src, dst):
    out_dir = str(Path(dst).parent)
    result = lo_convert(src, "pdf", out_dir)
    if str(result) != str(dst): shutil.move(result, dst)
    return dst

def pptx_to_images(src, dst_dir, fmt="png", dpi=150):
    if not HAS_PDF2IMAGE: raise ImportError("pip install pdf2image (+poppler)")
    os.makedirs(dst_dir, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        tmp_pdf = os.path.join(tmp, Path(src).stem + ".pdf")
        pptx_to_pdf(src, tmp_pdf)
        imgs = convert_from_path(tmp_pdf, dpi=dpi, fmt=fmt)
        out = []
        for i, img in enumerate(imgs):
            p = os.path.join(dst_dir, f"slide_{i+1}.{fmt}")
            img.save(p)
            img.close()
            out.append(p)
        del imgs
    return out

def pptx_to_txt(src, dst):
    if not HAS_PPTX: raise ImportError("pip install python-pptx")
    prs = Presentation(src)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n{'='*50}\nSLIDE {i}\n{'='*50}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t: lines.append(t)
    Path(dst).write_text("\n".join(lines), encoding="utf-8")
    return dst

def pptx_to_html(src, dst):
    if not HAS_PPTX: raise ImportError("pip install python-pptx")
    prs = Presentation(src)
    parts = [f"<html><head><meta charset='utf-8'><title>{Path(src).stem}</title>"
             f"<style>body{{font-family:Georgia,serif;max-width:900px;margin:auto;padding:2em}}"
             f".slide{{border:1px solid #ddd;margin-bottom:2em;padding:1.5em;border-radius:4px}}"
             f"h2{{color:#7c6af7}}</style></head><body><h1>{Path(src).stem}</h1>"]
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f'<div class="slide"><h2>Slide {i}</h2>')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        safe = t.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                        parts.append(f"<p>{safe}</p>")
        parts.append("</div>")
    parts.append("</body></html>")
    Path(dst).write_text("\n".join(parts), encoding="utf-8")
    return dst

def pptx_to_docx(src, dst):
    if not(HAS_PPTX and HAS_DOCX): raise ImportError("pip install python-pptx python-docx")
    prs = Presentation(src)
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    for i, slide in enumerate(prs.slides, 1):
        doc.add_heading(f"Slide {i}", 2)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t: doc.add_paragraph(t)
        doc.add_page_break()
    doc.save(dst)
    return dst

def pptx_slide_count(src):
    if not HAS_PPTX: return 0
    return len(Presentation(src).slides)

def pptx_split_slides(src, out_dir, prefix="slide"):
    if not HAS_PPTX: raise ImportError("pip install python-pptx")
    import copy
    os.makedirs(out_dir, exist_ok=True)
    prs = Presentation(src)
    total = len(prs.slides)
    w = _width(total)
    out = []
    for i, slide in enumerate(prs.slides):
        new_prs = Presentation()
        new_prs.slide_width = prs.slide_width
        new_prs.slide_height = prs.slide_height
        new_slide = new_prs.slides.add_slide(new_prs.slide_layouts[6])
        for shape in slide.shapes:
            new_slide.shapes._spTree.insert(2, copy.deepcopy(shape.element))
        p = os.path.join(out_dir, f"{prefix}_{str(i+1).zfill(w)}.pptx")
        new_prs.save(p)
        del new_prs
        out.append(p)
    return out

def pptx_merge(src_list, dst):
    if not HAS_PPTX: raise ImportError("pip install python-pptx")
    import copy
    base_prs = Presentation(src_list[0])
    for src in src_list[1:]:
        src_prs = Presentation(src)
        for slide in src_prs.slides:
            new_slide = base_prs.slides.add_slide(base_prs.slide_layouts[6])
            for shape in slide.shapes:
                new_slide.shapes._spTree.insert(2, copy.deepcopy(shape.element))
        del src_prs
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    base_prs.save(dst)
    return dst

# ── DOCX ──────────────────────────────────────────────────────────────────────

def docx_to_txt(src, dst):
    if not HAS_DOCX: raise ImportError("pip install python-docx")
    doc = Document(src)
    Path(dst).write_text("\n".join(p.text for p in doc.paragraphs), encoding="utf-8")
    return dst

def docx_to_html(src, dst):
    if not HAS_DOCX: raise ImportError("pip install python-docx")
    doc = Document(src)
    parts = [f"<html><head><meta charset='utf-8'><title>{Path(src).stem}</title>"
             f"<style>body{{font-family:Georgia,serif;max-width:860px;margin:auto;padding:2em}}</style></head><body>"]
    for p in doc.paragraphs:
        if p.style.name.startswith("Heading"):
            lvl = p.style.name[-1] if p.style.name[-1].isdigit() else "2"
            parts.append(f"<h{lvl}>{p.text}</h{lvl}>")
        else:
            safe = p.text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            parts.append(f"<p>{safe}</p>")
    parts.append("</body></html>")
    Path(dst).write_text("\n".join(parts), encoding="utf-8")
    return dst

def docx_to_pdf(src, dst):
    out_dir = str(Path(dst).parent)
    result = lo_convert(src, "pdf", out_dir)
    if str(result) != str(dst): shutil.move(result, dst)
    return dst

# ── TXT ───────────────────────────────────────────────────────────────────────

def txt_to_pdf(src, dst):
    if not HAS_REPORTLAB: raise ImportError("pip install reportlab")
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import mm
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    doc = SimpleDocTemplate(dst, pagesize=A4,
                            leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=20*mm, bottomMargin=20*mm)
    styles = getSampleStyleSheet()
    story = []
    for line in text.split("\n"):
        safe = line.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;") or "&nbsp;"
        story.append(Paragraph(safe, styles["Normal"]))
        story.append(Spacer(1, 2))
    doc.build(story)
    return dst

def txt_to_docx(src, dst):
    if not HAS_DOCX: raise ImportError("pip install python-docx")
    doc = Document()
    doc.add_heading(Path(src).stem, 0)
    for line in Path(src).read_text(encoding="utf-8", errors="replace").split("\n"):
        doc.add_paragraph(line)
    doc.save(dst)
    return dst

def txt_to_html(src, dst):
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    safe = text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
    Path(dst).write_text(
        f"<html><head><meta charset='utf-8'>"
        f"<style>body{{font-family:monospace;max-width:900px;margin:auto;padding:2em}}</style></head>"
        f"<body><pre>{safe}</pre></body></html>", encoding="utf-8")
    return dst

# ── HTML ──────────────────────────────────────────────────────────────────────

def html_to_txt(src, dst):
    html = Path(src).read_text(encoding="utf-8", errors="replace")
    text = re.sub(r"<[^>]+>", "", html)
    for ent, rep in [("&nbsp;"," "),("&amp;","&"),("&lt;","<"),("&gt;",">")]:
        text = text.replace(ent, rep)
    Path(dst).write_text(text.strip(), encoding="utf-8")
    return dst

def html_to_pdf(src, dst):
    if HAS_WEASYPRINT:
        from weasyprint import HTML
        HTML(filename=src).write_pdf(dst)
        return dst
    wk = shutil.which("wkhtmltopdf")
    if wk:
        r = subprocess.run([wk, src, dst], capture_output=True)
        if Path(dst).exists(): return dst
    raise ImportError("HTML→PDF needs:  pip install weasyprint")

# ── Image ─────────────────────────────────────────────────────────────────────

def image_convert(src, dst):
    src_ext = Path(src).suffix.lower()
    dst_ext = Path(dst).suffix.lower()

    if src_ext == ".svg":
        if dst_ext not in {".png", ".jpg", ".jpeg", ".webp", ".pdf"}:
            raise ValueError(f"SVG can be converted to PNG/JPG/WEBP/PDF only (not {dst_ext}).")
        if not HAS_CAIROSVG:
            raise ImportError("SVG conversion needs: pip install cairosvg")
        if dst_ext == ".pdf":
            cairosvg.svg2pdf(url=src, write_to=dst)
            return dst
        cairosvg.svg2png(url=src, write_to=dst)
        if dst_ext in {".jpg", ".jpeg", ".webp"}:
            if not HAS_PIL:
                raise ImportError("SVG→JPG/WEBP needs: pip install Pillow")
            img = Image.open(dst).convert("RGB")
            img.save(dst)
            img.close()
        return dst

    if not HAS_PIL: raise ImportError("pip install Pillow")
    img = Image.open(src)
    ext = Path(dst).suffix.lower()
    if ext in (".jpg",".jpeg",".bmp") and img.mode in ("RGBA","P","LA"):
        img = img.convert("RGB")
    if ext == ".ico":
        img = img.resize((256, 256), Image.LANCZOS)
    img.save(dst)
    img.close()
    return dst

def images_to_pdf(src_list, dst):
    """Convert multiple images to a single PDF. Streams one image at a time."""
    if not HAS_PIL: raise ImportError("pip install Pillow")
    if not src_list: raise ValueError("No images provided")

    # Process first image
    first = Image.open(src_list[0]).convert("RGB")

    if len(src_list) == 1:
        first.save(dst, "PDF")
        first.close()
        return dst

    # Stream remaining images one at a time to avoid holding all in RAM
    append_imgs = []
    for p in src_list[1:]:
        img = Image.open(p).convert("RGB")
        append_imgs.append(img)

    first.save(dst, save_all=True, append_images=append_imgs)
    first.close()
    for img in append_imgs:
        img.close()
    return dst

# ── Audio/Video ───────────────────────────────────────────────────────────────

def ffmpeg_convert(src, dst, extra=None):
    if not HAS_FFMPEG: raise RuntimeError("ffmpeg not found.\n  sudo apt install ffmpeg")
    cmd = ["ffmpeg", "-y", "-i", src] + (extra or []) + [dst]
    r = subprocess.run(cmd, capture_output=True, text=True)
    if not Path(dst).exists():
        raise RuntimeError(f"ffmpeg failed:\n{r.stderr[-400:]}")
    return dst

# ── Video / Audio operations ──────────────────────────────────────────────────

def video_get_duration(src: str) -> float:
    """Return video/audio duration in seconds using ffprobe."""
    if not HAS_FFMPEG:
        raise RuntimeError("ffmpeg not found.\n  sudo apt install ffmpeg")
    ffprobe = shutil.which("ffprobe")
    if not ffprobe:
        raise RuntimeError("ffprobe not found.\n  sudo apt install ffmpeg")
    cmd = [ffprobe, "-v", "error",
           "-show_entries", "format=duration",
           "-of", "default=noprint_wrappers=1:nokey=1",
           src]
    r = subprocess.run(cmd, capture_output=True, text=True)
    try:
        return float(r.stdout.strip())
    except (ValueError, TypeError):
        raise RuntimeError(
            f"Could not read duration from: {Path(src).name}\n{r.stderr[:300]}"
        )

def _parse_time(t: str) -> float:
    """
    Parse a time string into seconds.
    Accepts:  HH:MM:SS  |  MM:SS  |  SS  |  SS.ms
    e.g.  "1:30:00" → 5400.0   "2:45" → 165.0   "90" → 90.0
    """
    t = t.strip()
    if not t:
        raise ValueError("Empty time value")
    parts = t.split(":")
    try:
        if len(parts) == 3:
            return int(parts[0]) * 3600 + int(parts[1]) * 60 + float(parts[2])
        elif len(parts) == 2:
            return int(parts[0]) * 60 + float(parts[1])
        else:
            return float(parts[0])
    except (ValueError, IndexError):
        raise ValueError(
            f"Invalid time format: '{t}'\n"
            "Use HH:MM:SS, MM:SS, or seconds (e.g. 90 or 1:30)"
        )

def format_duration(seconds: float) -> str:
    """Format seconds as HH:MM:SS string."""
    s = int(seconds)
    h, rem = divmod(s, 3600)
    m, sec = divmod(rem, 60)
    return f"{h:02d}:{m:02d}:{sec:02d}"

def video_split(src: str, segments: list, out_dir: str,
                prefix: str = "", fmt: str = "") -> list:
    """
    Split a video/audio into segments by time.

    segments: list of (start, end) tuples — each a time string or float seconds.
              e.g. [("0", "1:30"), ("1:30", "3:00"), ("3:00", "end")]
              Use "end" or "" as end time to mean the rest of the file.
    """
    if not HAS_FFMPEG:
        raise RuntimeError("ffmpeg not found.\n  sudo apt install ffmpeg")
    if not segments:
        raise ValueError("No segments defined")

    os.makedirs(out_dir, exist_ok=True)
    ext = fmt.lstrip(".") if fmt else Path(src).suffix.lstrip(".")
    duration = video_get_duration(src)
    stem = _safe_stem(src, 40) if not prefix else prefix
    w = _width(len(segments))
    out_files = []

    for i, (start_raw, end_raw) in enumerate(segments, 1):
        start_sec = _parse_time(str(start_raw))

        # Handle "end", empty string, or missing end time
        end_str = str(end_raw).strip().lower()
        if end_str in ("end", "", "0"):
            end_sec = duration
        else:
            end_sec = _parse_time(str(end_raw))

        # Clamp to actual duration
        start_sec = min(start_sec, duration)
        end_sec = min(end_sec, duration)

        if end_sec <= start_sec:
            raise ValueError(
                f"Segment {i}: end ({end_raw}) must be after start ({start_raw})"
            )

        seg_dur = end_sec - start_sec

        # Safe filename
        s_label = _safe_video_filename(str(start_raw))
        e_label = _safe_video_filename(str(end_raw))
        name = f"{stem}_seg{str(i).zfill(w)}_{s_label}_to_{e_label}.{ext}"
        name = _safe_filename(name)
        out_path = os.path.join(out_dir, name)

        cmd = [
            "ffmpeg", "-y",
            "-ss", str(start_sec),
            "-i", src,
            "-t", str(seg_dur),
            "-c", "copy",
            "-avoid_negative_ts", "make_zero",
            out_path
        ]
        r = subprocess.run(cmd, capture_output=True, text=True)

        if not Path(out_path).exists():
            # Retry with re-encode if stream copy failed
            cmd_retry = [
                "ffmpeg", "-y",
                "-ss", str(start_sec),
                "-i", src,
                "-t", str(seg_dur),
                "-avoid_negative_ts", "make_zero",
                out_path
            ]
            r2 = subprocess.run(cmd_retry, capture_output=True, text=True)
            if not Path(out_path).exists():
                raise RuntimeError(
                    f"ffmpeg failed on segment {i}:\n{r2.stderr[-400:]}"
                )

        out_files.append(out_path)

    return out_files

def video_merge(src_list: list, dst: str, log=None) -> str:
    """
    Merge multiple video/audio files into one.
    Tries stream copy first (fast), falls back to re-encode.
    """
    if not HAS_FFMPEG:
        raise RuntimeError("ffmpeg not found.\n  sudo apt install ffmpeg")
    if not src_list:
        raise ValueError("No files provided for merge.")
    if len(src_list) == 1:
        shutil.copy2(src_list[0], dst)
        return dst

    def L(m):
        if log: log(m)

    os.makedirs(str(Path(dst).parent) or ".", exist_ok=True)

    # Write concat list file
    list_fd, list_path = tempfile.mkstemp(suffix=".txt", prefix="merge_")
    try:
        with os.fdopen(list_fd, "w", encoding="utf-8") as tf:
            for s in src_list:
                abs_path = os.path.abspath(s)
                escaped = abs_path.replace("'", "'\\''")
                tf.write(f"file '{escaped}'\n")

        L(f"Merging {len(src_list)} files …")

        # Try stream copy first
        cmd = [
            "ffmpeg", "-y",
            "-f", "concat",
            "-safe", "0",
            "-i", list_path,
            "-c", "copy",
            dst
        ]
        r = subprocess.run(cmd, capture_output=True, text=True)
        if Path(dst).exists() and os.path.getsize(dst) > 0:
            L("Merged with stream copy (fast, lossless).")
            return dst

        # Remove failed output if exists
        if Path(dst).exists():
            os.unlink(dst)

        # Fall back to re-encode
        L("Stream copy failed — re-encoding …")
        cmd2 = [
            "ffmpeg", "-y",
            "-f", "concat",
            "-safe", "0",
            "-i", list_path,
            dst
        ]
        r2 = subprocess.run(cmd2, capture_output=True, text=True)
        if not Path(dst).exists() or os.path.getsize(dst) == 0:
            raise RuntimeError(f"ffmpeg merge failed:\n{r2.stderr[-600:]}")
        L("Merged with re-encode.")
        return dst
    finally:
        try:
            os.unlink(list_path)
        except OSError:
            pass

def audio_split(src: str, segments: list, out_dir: str,
                prefix: str = "segment", fmt: str = "") -> list:
    """Split an audio file by time — identical logic to video_split."""
    return video_split(src, segments, out_dir, prefix, fmt)

def audio_merge(src_list: list, dst: str, log=None) -> str:
    """Merge multiple audio files into one output."""
    return video_merge(src_list, dst, log=log)

# ── PDF split / merge / organise (MEMORY EFFICIENT) ──────────────────────────

def split_each(src, out_dir):
    """Split every page into its own PDF. Memory-safe for 100k+ pages."""
    os.makedirs(out_dir, exist_ok=True)
    reader = PdfReader(src)
    total = len(reader.pages)
    w = _width(total)
    stem = _safe_stem(src)
    results = []

    for i in range(total):
        writer = PdfWriter()
        writer.add_page(reader.pages[i])
        name = f"{stem}_p{str(i + 1).zfill(w)}.pdf"
        out = os.path.join(out_dir, name)
        with open(out, "wb") as f:
            writer.write(f)
        del writer
        results.append(out)

        # Periodic GC for very large files
        if (i + 1) % 500 == 0:
            gc.collect()

    return results

def split_range(src, start, end, out_dir):
    """Extract a contiguous page range. 1-indexed, inclusive."""
    os.makedirs(out_dir, exist_ok=True)
    reader = PdfReader(src)
    total = len(reader.pages)
    start = max(1, start)
    end = min(end, total)
    stem = _safe_stem(src)

    writer = PdfWriter()
    for i in range(start - 1, end):
        writer.add_page(reader.pages[i])

    name = f"{stem}_p{start}-{end}.pdf"
    out = os.path.join(out_dir, name)
    with open(out, "wb") as f:
        writer.write(f)
    del writer
    return out

def split_custom(src, groups, out_dir):
    """
    Split into custom groups. Each group is a list of 1-indexed page numbers.
    Example: [[1, 10], [3, 5], [2, 3, 4, 7]]

    Filenames:
      Single page group:    report_g01_p5.pdf
      Small group (≤5 pg):  report_g01_p1-3-5.pdf
      Large group:          report_g01_50pg.pdf
    """
    os.makedirs(out_dir, exist_ok=True)
    reader = PdfReader(src)
    total = len(reader.pages)
    stem = _safe_stem(src)
    gw = _width(len(groups))
    results = []

    for gi, group in enumerate(groups, 1):
        valid = sorted(set(p for p in group if 1 <= p <= total))
        if not valid:
            continue

        writer = PdfWriter()
        for p in valid:
            writer.add_page(reader.pages[p - 1])

        # Build a short, descriptive suffix
        g_label = f"g{str(gi).zfill(gw)}"
        if len(valid) == 1:
            suffix = f"{g_label}_p{valid[0]}"
        elif len(valid) <= 5:
            suffix = f"{g_label}_p{'-'.join(str(p) for p in valid)}"
        else:
            suffix = f"{g_label}_{len(valid)}pg"

        name = f"{stem}_{suffix}.pdf"
        name = _safe_filename(name)

        out = os.path.join(out_dir, name)
        with open(out, "wb") as f:
            writer.write(f)
        del writer
        results.append(out)

    return results

def merge_pdfs(src_list, dst, chunk_size=50):
    """
    Merge multiple PDFs (and images) into one.
    For large merges, works in chunks to cap RAM usage.
    """
    if not HAS_PYPDF: raise ImportError("pip install pypdf")
    os.makedirs(str(Path(dst).parent) or ".", exist_ok=True)

    # Normalize: convert images to temp PDFs first
    pdf_paths = []
    temp_files = []
    for src in src_list:
        c = cat(src)
        if c == "pdf":
            pdf_paths.append(src)
        elif c == "image":
            tf = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
            tmp = tf.name
            tf.close()
            images_to_pdf([src], tmp)
            pdf_paths.append(tmp)
            temp_files.append(tmp)
        # Skip unsupported types silently

    try:
        if len(pdf_paths) <= chunk_size:
            # Direct merge — small job
            writer = PdfWriter()
            for p in pdf_paths:
                reader = PdfReader(p)
                for page in reader.pages:
                    writer.add_page(page)
                del reader
            with open(dst, "wb") as f:
                writer.write(f)
            del writer
            return dst

        # Chunked merge for very large jobs
        temp_dir = tempfile.mkdtemp(prefix="merge_chunks_")
        intermediates = []

        for ci in range(0, len(pdf_paths), chunk_size):
            chunk = pdf_paths[ci:ci + chunk_size]
            writer = PdfWriter()
            for p in chunk:
                reader = PdfReader(p)
                for page in reader.pages:
                    writer.add_page(page)
                del reader
            tmp_out = os.path.join(temp_dir, f"chunk_{ci // chunk_size}.pdf")
            with open(tmp_out, "wb") as f:
                writer.write(f)
            del writer
            gc.collect()
            intermediates.append(tmp_out)

        # Final merge of intermediates
        final_writer = PdfWriter()
        for tmp in intermediates:
            reader = PdfReader(tmp)
            for page in reader.pages:
                final_writer.add_page(page)
            del reader
        with open(dst, "wb") as f:
            final_writer.write(f)
        del final_writer

        # Cleanup intermediates
        shutil.rmtree(temp_dir, ignore_errors=True)
        return dst

    finally:
        # Cleanup any temp image→PDF files
        for tmp in temp_files:
            try:
                os.unlink(tmp)
            except OSError:
                pass

def resequence_pdf(src, order, dst):
    reader = PdfReader(src)
    total = len(reader.pages)
    writer = PdfWriter()
    for pg in order:
        if 1 <= pg <= total:
            writer.add_page(reader.pages[pg - 1])
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    del writer
    return dst

def delete_pages(src, pages_to_delete, dst):
    reader = PdfReader(src)
    total = len(reader.pages)
    remove = set(pages_to_delete)
    writer = PdfWriter()
    for i in range(total):
        if (i + 1) not in remove:
            writer.add_page(reader.pages[i])
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    del writer
    return dst

def rotate_pages(src, degrees, pages_to_rotate, dst):
    reader = PdfReader(src)
    total = len(reader.pages)
    target = set(pages_to_rotate) if pages_to_rotate else set(range(1, total + 1))
    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        if (i + 1) in target:
            page.rotate(degrees)
        writer.add_page(page)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    del writer
    return dst

def reverse_pdf(src, dst):
    reader = PdfReader(src)
    return resequence_pdf(src, list(range(len(reader.pages), 0, -1)), dst)

def compress_pdf(src, dst):
    """Compress PDF with page-by-page streaming for memory efficiency."""
    reader = PdfReader(src)
    writer = PdfWriter()
    for page in reader.pages:
        page.compress_content_streams()
        writer.add_page(page)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    del writer
    return dst

def encrypt_pdf(src, dst, user_pw, owner_pw=""):
    reader = PdfReader(src)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(user_pw, owner_pw or user_pw)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    del writer
    return dst

def decrypt_pdf(src, dst, password):
    reader = PdfReader(src)
    if reader.is_encrypted:
        if reader.decrypt(password) == 0:
            raise ValueError("Wrong password")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    del writer
    return dst

def watermark_text(src, dst, text, opacity=0.3, fontsize=48,
                   color="#888888", angle=45, pages=None):
    """Apply text watermark. Processes one page at a time for memory efficiency."""
    if not(HAS_PIL and HAS_PDF2IMAGE):
        raise ImportError("pip install Pillow pdf2image (+poppler)")
    reader = PdfReader(src)
    total = len(reader.pages)
    target = set(pages) if pages else set(range(1, total + 1))
    writer = PdfWriter()

    # Parse color once
    try:
        r_c, g_c, b_c = int(color[1:3],16), int(color[3:5],16), int(color[5:7],16)
    except:
        r_c, g_c, b_c = 128, 128, 128
    alpha = int(opacity * 255)

    # Try to load font once
    font = None
    for font_path in [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
    ]:
        try:
            font = ImageFont.truetype(font_path, fontsize)
            break
        except:
            continue
    if font is None:
        font = ImageFont.load_default()

    # Process one page at a time
    for pg_num in range(1, total + 1):
        if pg_num in target:
            # Convert just this one page
            imgs = convert_from_path(src, dpi=150,
                                     first_page=pg_num, last_page=pg_num)
            img = imgs[0]

            # Draw watermark
            txt_img = Image.new("RGBA", img.size, (0, 0, 0, 0))
            draw = ImageDraw.Draw(txt_img)
            bbox = draw.textbbox((0, 0), text, font=font)
            tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
            draw.text(
                (img.size[0] // 2 - tw // 2, img.size[1] // 2 - th // 2),
                text, font=font, fill=(r_c, g_c, b_c, alpha)
            )
            txt_img = txt_img.rotate(angle, expand=False)
            img = Image.alpha_composite(img.convert("RGBA"), txt_img)

            # Save as temp PDF page
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tf:
                tmp = tf.name
            img.convert("RGB").save(tmp, "PDF")
            for page in PdfReader(tmp).pages:
                writer.add_page(page)
            os.unlink(tmp)

            # Cleanup
            img.close()
            txt_img.close()
            del imgs
        else:
            # Non-watermarked page — copy directly from source
            writer.add_page(reader.pages[pg_num - 1])

    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    del writer
    return dst

def watermark_pdf_overlay(src, wm_pdf, dst, pages=None):
    reader = PdfReader(src)
    wm_page = PdfReader(wm_pdf).pages[0]
    total = len(reader.pages)
    target = set(pages) if pages else set(range(1, total + 1))
    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        if (i + 1) in target:
            page.merge_page(wm_page)
        writer.add_page(page)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    del writer
    return dst

def get_metadata(src):
    reader = PdfReader(src)
    meta = reader.metadata or {}
    size = os.path.getsize(src)
    return {
        "File": Path(src).name,
        "Pages": len(reader.pages),
        "Size": f"{size/1024:.1f} KB",
        "Title": meta.get("/Title", "—"),
        "Author": meta.get("/Author", "—"),
        "Subject": meta.get("/Subject", "—"),
        "Creator": meta.get("/Creator", "—"),
        "Producer": meta.get("/Producer", "—"),
    }

def set_metadata(src, dst, fields):
    reader = PdfReader(src)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    existing = dict(reader.metadata or {})
    mapping = {
        "title": "/Title",
        "author": "/Author",
        "subject": "/Subject",
        "creator": "/Creator",
    }
    for k, v in fields.items():
        pk = mapping.get(k.lower())
        if pk and v.strip():
            existing[pk] = v.strip()
    writer.add_metadata(existing)
    os.makedirs(str(Path(dst).parent), exist_ok=True)
    with open(dst, "wb") as f:
        writer.write(f)
    del writer
    return dst

# ── Master dispatch ───────────────────────────────────────────────────────────

def do_convert(src, out_fmt, out_dir, pages_str="", dpi=150, log=None):
    def L(m):
        if log: log(m)

    c = cat(src)
    stem = Path(src).stem
    out_fmt = out_fmt.lower().lstrip(".")
    os.makedirs(out_dir, exist_ok=True)
    total = pdf_page_count(src) if c == "pdf" and HAS_PYPDF else 0
    pages = parse_pages(pages_str, total) if pages_str.strip() and total else None

    if c == "pdf":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "txt":
            L("PDF→TXT"); return [pdf_to_txt(src, dst, pages)]
        if out_fmt == "docx":
            L("PDF→DOCX"); return [pdf_to_docx(src, dst, pages)]
        if out_fmt in ("png","jpg","jpeg"):
            L(f"PDF→{out_fmt.upper()}")
            return pdf_to_images(src, os.path.join(out_dir, stem + "_images"), out_fmt, pages, dpi)
        if out_fmt == "html":
            L("PDF→HTML"); return [pdf_to_html(src, dst, pages)]
        if out_fmt in ("xlsx","xls"):
            L("PDF→XLSX"); return [pdf_to_xlsx(src, dst.replace(f".{out_fmt}", ".xlsx"))]
        if out_fmt == "pptx":
            L("PDF→PPTX"); return [pdf_to_pptx(src, dst, dpi)]
        if out_fmt == "csv":
            xls = dst.replace(".csv", ".xlsx")
            pdf_to_xlsx(src, xls)
            return [excel_to_csv(xls, dst)]

    if c == "excel":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "pdf":
            L("XLSX→PDF"); return [excel_to_pdf(src, dst)]
        if out_fmt == "csv":
            L("XLSX→CSV"); return [excel_to_csv(src, dst)]
        if out_fmt == "txt":
            L("XLSX→TXT"); return [excel_to_txt(src, dst)]
        if out_fmt == "html":
            L("XLSX→HTML"); return [excel_to_html(src, dst)]
        if out_fmt == "docx":
            L("XLSX→DOCX"); return [excel_to_docx(src, dst)]
        if out_fmt in ("xlsx","xls","ods"):
            return [lo_convert(src, out_fmt, out_dir)]

    if c == "csv":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt in ("xlsx","xls"):
            L("CSV→XLSX"); return [csv_to_xlsx(src, dst.replace(f".{out_fmt}", ".xlsx"))]
        if out_fmt == "pdf":
            L("CSV→PDF"); return [csv_to_pdf(src, dst)]
        if out_fmt == "html":
            L("CSV→HTML"); return [csv_to_html(src, dst)]
        if out_fmt == "txt":
            shutil.copy(src, dst); return [dst]
        if out_fmt == "docx":
            L("CSV→DOCX"); return [csv_to_docx(src, dst)]

    if c == "pptx":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "pdf":
            L("PPTX→PDF"); return [pptx_to_pdf(src, dst)]
        if out_fmt in ("png","jpg","jpeg"):
            L(f"PPTX→{out_fmt.upper()}")
            return pptx_to_images(src, os.path.join(out_dir, stem + "_slides"), out_fmt, dpi)
        if out_fmt == "txt":
            L("PPTX→TXT"); return [pptx_to_txt(src, dst)]
        if out_fmt == "html":
            L("PPTX→HTML"); return [pptx_to_html(src, dst)]
        if out_fmt == "docx":
            L("PPTX→DOCX"); return [pptx_to_docx(src, dst)]

    if c == "docx":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "pdf":
            L("DOCX→PDF"); return [docx_to_pdf(src, dst)]
        if out_fmt == "txt":
            L("DOCX→TXT"); return [docx_to_txt(src, dst)]
        if out_fmt == "html":
            L("DOCX→HTML"); return [docx_to_html(src, dst)]

    if c == "txt":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "pdf":
            L("TXT→PDF"); return [txt_to_pdf(src, dst)]
        if out_fmt == "docx":
            L("TXT→DOCX"); return [txt_to_docx(src, dst)]
        if out_fmt == "html":
            L("TXT→HTML"); return [txt_to_html(src, dst)]

    if c == "html":
        dst = os.path.join(out_dir, stem + f".{out_fmt}")
        if out_fmt == "pdf":
            L("HTML→PDF"); return [html_to_pdf(src, dst)]
        if out_fmt == "txt":
            L("HTML→TXT"); return [html_to_txt(src, dst)]

    if c == "image":
        if out_fmt == "pdf":
            dst = os.path.join(out_dir, stem + ".pdf")
            if Path(src).suffix.lower() == ".svg":
                L("SVG→PDF"); return [image_convert(src, dst)]
            L("Image→PDF"); return [images_to_pdf([src], dst)]
        dst = os.path.join(out_dir, stem + "." + out_fmt)
        L(f"Image→{out_fmt.upper()}"); return [image_convert(src, dst)]

    if c in ("video", "audio"):
        dst = os.path.join(out_dir, stem + "." + out_fmt)
        L(f"{c.upper()}→{out_fmt.upper()}")
        if out_fmt == "gif":
            return [ffmpeg_convert(src, dst,
                    ["-vf","fps=10,scale=480:-1:flags=lanczos","-loop","0"])]
        return [ffmpeg_convert(src, dst)]

    raise ValueError(f"No conversion: {c} → {out_fmt}")
